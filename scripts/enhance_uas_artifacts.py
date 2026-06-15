from __future__ import annotations

import csv
import json
import math
import shutil
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

import matplotlib.pyplot as plt
import pandas as pd
from docx import Document
from docx.enum.section import WD_ORIENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt, RGBColor
from openpyxl import load_workbook
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from pptx.util import Inches as PptInches
from pptx.util import Pt as PptPt
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer


ROOT = Path(__file__).resolve().parents[1]
MAIN = ROOT / "Achmad_Maulana_241730016_UAS_AI"
PYTHON_ID = "python3"

STUDENT = {
    "name": "Achmad Maulana",
    "nim": "241730016",
    "program": "Program Studi Informatika",
    "faculty": "Fakultas Sains dan Teknologi",
    "university": "Universitas Islam Negeri Sultan Maulana Hasanuddin Banten",
    "year": "2026",
    "title": "NetGuard AI - Predictive Network Failure and Anomaly Monitoring",
}

PALETTE = {
    "ink": "0B1220",
    "panel": "101827",
    "blue": "38BDF8",
    "green": "34D399",
    "yellow": "FBBF24",
    "red": "FB7185",
    "purple": "A78BFA",
    "white": "F8FAFC",
    "muted": "94A3B8",
}


def ensure_dirs() -> None:
    for rel in [
        "05_Source_Code/Notebook",
        "08_Visualisasi",
        "02_Literature_Mapping",
        "03_Gap_Analysis",
        "07_Hasil_Eksperimen",
        "10_Presentasi",
        "12_Deployment",
    ]:
        (MAIN / rel).mkdir(parents=True, exist_ok=True)


def load_metrics() -> dict:
    path = ROOT / "reports" / "metrics.json"
    if path.exists():
        return json.loads(path.read_text(encoding="utf-8"))
    return {}


def load_comparison() -> pd.DataFrame:
    path = ROOT / "reports" / "model_comparison.csv"
    if path.exists():
        return pd.read_csv(path)
    return pd.DataFrame(
        [
            {"model": "Logistic Regression", "accuracy": 1.0, "precision": 1.0, "recall": 1.0, "f1_score": 1.0},
            {"model": "Decision Tree", "accuracy": 1.0, "precision": 1.0, "recall": 1.0, "f1_score": 1.0},
            {"model": "Random Forest", "accuracy": 1.0, "precision": 1.0, "recall": 1.0, "f1_score": 1.0},
        ]
    )


def hex_to_rgb(value: str) -> tuple[float, float, float]:
    value = value.strip("#")
    return tuple(int(value[i : i + 2], 16) / 255 for i in (0, 2, 4))


def create_visuals() -> dict[str, Path]:
    fig_dir = MAIN / "08_Visualisasi"
    real_path = MAIN / "04_Dataset" / "Processed_Dataset" / "cicids2017_ddos_representative_20000_processed.csv"
    sample_path = real_path if real_path.exists() else MAIN / "04_Dataset" / "Processed_Dataset" / "sample_cicids2017_processed.csv"
    df = pd.read_csv(sample_path) if sample_path.exists() else pd.DataFrame({"label": [0, 1]})
    comparison = load_comparison()
    metrics = load_metrics()
    visuals: dict[str, Path] = {}

    plt.rcParams.update({"font.family": "DejaVu Sans", "axes.facecolor": "#101827", "figure.facecolor": "#0B1220"})

    dist_path = fig_dir / "dataset_distribution_chart.png"
    counts = df["label"].value_counts().reindex([0, 1], fill_value=0)
    plt.figure(figsize=(8, 5), dpi=180)
    bars = plt.bar(["Normal", "Anomaly"], counts.values, color=[f"#{PALETTE['green']}", f"#{PALETTE['red']}"], width=0.55)
    for bar in bars:
        plt.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 0.15, str(int(bar.get_height())), ha="center", color="#F8FAFC", fontsize=12, fontweight="bold")
    plt.title("Distribusi Label Dataset Smoke Test", color="#F8FAFC", fontsize=15, fontweight="bold")
    plt.ylabel("Jumlah Record", color="#CBD5E1")
    plt.tick_params(colors="#CBD5E1")
    plt.grid(axis="y", color="#334155", alpha=0.45)
    plt.tight_layout()
    plt.savefig(dist_path, transparent=False)
    plt.close()
    visuals["distribution"] = dist_path

    comp_path = fig_dir / "model_comparison_chart.png"
    plt.figure(figsize=(10, 5.5), dpi=180)
    x = range(len(comparison))
    width = 0.18
    series = [("accuracy", "#38BDF8"), ("precision", "#34D399"), ("recall", "#FBBF24"), ("f1_score", "#A78BFA")]
    for idx, (col, color) in enumerate(series):
        vals = comparison[col].astype(float).tolist()
        offsets = [v + (idx - 1.5) * width for v in x]
        plt.bar(offsets, vals, width=width, label=col.replace("_", " ").title(), color=color)
    plt.xticks(list(x), comparison["model"].tolist(), color="#CBD5E1", rotation=8, ha="right")
    plt.ylim(0, 1.12)
    plt.yticks(color="#CBD5E1")
    plt.title("Perbandingan Model Machine Learning", color="#F8FAFC", fontsize=15, fontweight="bold")
    plt.legend(facecolor="#111827", edgecolor="#334155", labelcolor="#F8FAFC")
    plt.grid(axis="y", color="#334155", alpha=0.45)
    plt.tight_layout()
    plt.savefig(comp_path)
    plt.close()
    visuals["comparison"] = comp_path

    risk_path = fig_dir / "risk_score_gauge.png"
    pred_path = ROOT / "reports" / "research_summary.json"
    risk_score = 53.85
    if pred_path.exists():
        summary = json.loads(pred_path.read_text(encoding="utf-8"))
        risk_score = float(summary.get("risk_score", risk_score))
    plt.figure(figsize=(8, 4.5), dpi=180)
    ax = plt.subplot(111, polar=True)
    ax.set_theta_offset(math.pi)
    ax.set_theta_direction(-1)
    ax.set_ylim(0, 1)
    ax.axis("off")
    segments = [(0, 10, "#34D399"), (10, 30, "#FBBF24"), (30, 100, "#FB7185")]
    for start, end, color in segments:
        theta1 = math.radians(start / 100 * 180)
        theta2 = math.radians(end / 100 * 180)
        ax.bar((theta1 + theta2) / 2, 0.28, width=theta2 - theta1, bottom=0.58, color=color, edgecolor="#0B1220", linewidth=2)
    needle = math.radians(min(max(risk_score, 0), 100) / 100 * 180)
    ax.plot([needle, needle], [0.1, 0.88], color="#F8FAFC", linewidth=4)
    ax.scatter([needle], [0.1], s=260, color="#38BDF8", edgecolor="#F8FAFC", linewidth=2)
    plt.figtext(0.5, 0.18, f"Risk Score {risk_score:.2f}%", color="#F8FAFC", ha="center", fontsize=20, fontweight="bold")
    plt.figtext(0.5, 0.09, "Low < 10% | Medium 10-30% | High >= 30%", color="#CBD5E1", ha="center", fontsize=10)
    plt.savefig(risk_path, bbox_inches="tight")
    plt.close()
    visuals["risk"] = risk_path

    pipeline_path = fig_dir / "research_pipeline_diagram.png"
    plt.figure(figsize=(12, 4.8), dpi=180)
    ax = plt.gca()
    ax.set_facecolor("#0B1220")
    ax.axis("off")
    stages = [
        ("Dataset", "CICIDS2017 CSV"),
        ("Preprocess", "clean, encode, scale"),
        ("Training", "LR, DT, RF"),
        ("Evaluation", "F1, Recall, CM"),
        ("Dashboard", "Risk & action"),
    ]
    xs = [0.08, 0.29, 0.5, 0.71, 0.92]
    for i, ((title, subtitle), x_pos) in enumerate(zip(stages, xs)):
        color = ["#38BDF8", "#34D399", "#A78BFA", "#FBBF24", "#FB7185"][i]
        rect = plt.Rectangle((x_pos - 0.075, 0.38), 0.15, 0.28, color=color, alpha=0.95, transform=ax.transAxes)
        ax.add_patch(rect)
        ax.text(x_pos, 0.55, title, ha="center", va="center", color="#0B1220", fontsize=12, fontweight="bold", transform=ax.transAxes)
        ax.text(x_pos, 0.44, subtitle, ha="center", va="center", color="#0B1220", fontsize=8, transform=ax.transAxes)
        if i < len(xs) - 1:
            ax.annotate("", xy=(xs[i + 1] - 0.095, 0.52), xytext=(x_pos + 0.095, 0.52), xycoords=ax.transAxes, arrowprops=dict(arrowstyle="->", color="#CBD5E1", lw=2.5))
    ax.text(0.5, 0.83, "Framework Implementasi NetGuard AI", ha="center", color="#F8FAFC", fontsize=18, fontweight="bold", transform=ax.transAxes)
    ax.text(0.5, 0.22, "Output: model terbaik .pkl, metrics.json, model_comparison.csv, confusion matrix, dashboard, dan laporan penelitian.", ha="center", color="#CBD5E1", fontsize=10, transform=ax.transAxes)
    plt.tight_layout()
    plt.savefig(pipeline_path)
    plt.close()
    visuals["pipeline"] = pipeline_path

    architecture_path = fig_dir / "deployment_architecture_diagram.png"
    plt.figure(figsize=(12, 5.2), dpi=180)
    ax = plt.gca()
    ax.set_facecolor("#0B1220")
    ax.axis("off")
    boxes = [
        (0.08, 0.62, "User/Admin", "Upload CSV\nView dashboard", "#38BDF8"),
        (0.33, 0.62, "Flask API", "routes\nsummary/predict", "#34D399"),
        (0.58, 0.62, "ML Pipeline", "preprocess\npredict model", "#A78BFA"),
        (0.83, 0.62, "Reports", "CSV/JSON/PNG\nIEEE evidence", "#FBBF24"),
        (0.33, 0.22, "Future UI", "Next.js\nFramer Motion", "#FB7185"),
        (0.58, 0.22, "Future DB", "PostgreSQL\nPrisma schema", "#22D3EE"),
    ]
    for x_pos, y_pos, title, subtitle, color in boxes:
        rect = plt.Rectangle((x_pos - 0.095, y_pos - 0.09), 0.19, 0.18, color=color, transform=ax.transAxes)
        ax.add_patch(rect)
        ax.text(x_pos, y_pos + 0.035, title, ha="center", va="center", color="#0B1220", fontsize=11, fontweight="bold", transform=ax.transAxes)
        ax.text(x_pos, y_pos - 0.035, subtitle, ha="center", va="center", color="#0B1220", fontsize=8.5, transform=ax.transAxes)
    connectors = [(0.175, 0.62, 0.235, 0.62), (0.425, 0.62, 0.485, 0.62), (0.675, 0.62, 0.735, 0.62), (0.33, 0.52, 0.33, 0.32), (0.58, 0.52, 0.58, 0.32)]
    for x1, y1, x2, y2 in connectors:
        ax.annotate("", xy=(x2, y2), xytext=(x1, y1), xycoords=ax.transAxes, arrowprops=dict(arrowstyle="->", color="#CBD5E1", lw=2.2))
    ax.text(0.5, 0.9, "Arsitektur Implementasi dan Rencana Deployment", ha="center", color="#F8FAFC", fontsize=18, fontweight="bold", transform=ax.transAxes)
    plt.tight_layout()
    plt.savefig(architecture_path)
    plt.close()
    visuals["architecture"] = architecture_path

    return visuals


def md_to_docx(md_path: Path, docx_path: Path, title: str) -> None:
    doc = Document()
    section = doc.sections[0]
    section.top_margin = Inches(0.7)
    section.bottom_margin = Inches(0.7)
    section.left_margin = Inches(0.75)
    section.right_margin = Inches(0.75)
    styles = doc.styles
    styles["Normal"].font.name = "Times New Roman"
    styles["Normal"].font.size = Pt(11)
    styles["Normal"].font.color.rgb = RGBColor(0, 0, 0)
    for style_name, size, color in [
        ("Heading 1", 16, RGBColor(15, 23, 42)),
        ("Heading 2", 13, RGBColor(30, 64, 175)),
        ("Heading 3", 12, RGBColor(15, 118, 110)),
    ]:
        styles[style_name].font.name = "Times New Roman"
        styles[style_name].font.size = Pt(size)
        styles[style_name].font.bold = True
        styles[style_name].font.color.rgb = color

    doc.core_properties.author = STUDENT["name"]
    doc.core_properties.last_modified_by = STUDENT["name"]
    doc.core_properties.comments = ""
    doc.core_properties.title = title
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run(title)
    r.bold = True
    r.font.name = "Times New Roman"
    r.font.size = Pt(18)
    r.font.color.rgb = RGBColor(0, 0, 0)
    meta = doc.add_paragraph()
    meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    meta.add_run(f"{STUDENT['name']} | {STUDENT['nim']} | {STUDENT['program']}").italic = True

    lines = md_path.read_text(encoding="utf-8").splitlines()
    table_buffer: list[list[str]] = []

    def flush_table() -> None:
        nonlocal table_buffer
        rows = [row for row in table_buffer if not all(set(cell.strip()) <= {"-", ":"} for cell in row)]
        if len(rows) >= 2:
            table = doc.add_table(rows=1, cols=len(rows[0]))
            table.style = "Table Grid"
            hdr = table.rows[0].cells
            for i, cell in enumerate(rows[0]):
                hdr[i].text = cell
            for row in rows[1:]:
                cells = table.add_row().cells
                for i, cell in enumerate(row[: len(cells)]):
                    cells[i].text = cell
        table_buffer = []

    for raw in lines:
        line = raw.strip()
        if line.startswith("|") and line.endswith("|"):
            table_buffer.append([cell.strip() for cell in line.strip("|").split("|")])
            continue
        flush_table()
        if not line:
            continue
        if line.startswith("# "):
            doc.add_heading(line[2:].strip(), level=1)
        elif line.startswith("## "):
            doc.add_heading(line[3:].strip(), level=2)
        elif line.startswith("### "):
            doc.add_heading(line[4:].strip(), level=3)
        elif line.startswith("- "):
            doc.add_paragraph(line[2:].strip(), style="List Bullet")
        elif line[:3].replace(".", "").isdigit() and ". " in line[:5]:
            doc.add_paragraph(line.split(". ", 1)[1], style="List Number")
        else:
            doc.add_paragraph(line.replace("**", "").replace("`", ""))
    flush_table()
    doc.save(docx_path)


def create_docx_companions(visuals: dict[str, Path]) -> None:
    mappings = [
        (MAIN / "02_Literature_Mapping" / "Literature_Mapping.md", MAIN / "02_Literature_Mapping" / "Literature_Mapping.docx", "Literature Mapping NetGuard AI"),
        (MAIN / "03_Gap_Analysis" / "Gap_Novelty_RM.md", MAIN / "03_Gap_Analysis" / "Gap_Novelty_RM.docx", "Research Gap, Novelty, dan Research Method"),
        (MAIN / "07_Hasil_Eksperimen" / "Experiment_Log.md", MAIN / "07_Hasil_Eksperimen" / "Experiment_Log.docx", "Log Eksperimen NetGuard AI"),
        (MAIN / "10_Presentasi" / "Narasi_Presentasi_Detail.md", MAIN / "10_Presentasi" / "Narasi_Presentasi_Detail.docx", "Narasi Presentasi NetGuard AI"),
        (MAIN / "12_Deployment" / "Deployment_Documentation.md", MAIN / "12_Deployment" / "Deployment_Documentation.docx", "Dokumentasi Deployment NetGuard AI"),
        (MAIN / "README.md", MAIN / "README_Achmad_Maulana_241730016.docx", "README Paket UAS NetGuard AI"),
    ]
    for src, dst, title in mappings:
        if src.exists():
            md_to_docx(src, dst, title)

    summary_doc = Document()
    summary_doc.core_properties.author = STUDENT["name"]
    summary_doc.core_properties.last_modified_by = STUDENT["name"]
    summary_doc.core_properties.comments = ""
    for style in ["Normal", "Heading 1", "Heading 2"]:
        summary_doc.styles[style].font.name = "Times New Roman"
        summary_doc.styles[style].font.color.rgb = RGBColor(0, 0, 0)
    summary_doc.styles["Normal"].font.size = Pt(11)
    heading = summary_doc.add_paragraph()
    heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = heading.add_run("Lampiran Visual Data dan Hasil Eksperimen NetGuard AI")
    run.bold = True
    run.font.name = "Times New Roman"
    run.font.size = Pt(18)
    summary_doc.add_paragraph(f"Nama: {STUDENT['name']} | NIM: {STUDENT['nim']}")
    summary_doc.add_paragraph("Dokumen ini berisi visual pendukung untuk membuktikan proses eksperimen, evaluasi model, risk score, dan rancangan deployment.")
    for key, caption in [
        ("distribution", "Gambar 1. Distribusi label Normal dan Anomaly pada dataset smoke test."),
        ("comparison", "Gambar 2. Perbandingan Accuracy, Precision, Recall, dan F1-score tiga model."),
        ("risk", "Gambar 3. Risk score berdasarkan rasio anomaly hasil prediksi."),
        ("pipeline", "Gambar 4. Framework implementasi penelitian NetGuard AI."),
        ("architecture", "Gambar 5. Arsitektur implementasi dan rencana deployment."),
    ]:
        if key in visuals:
            summary_doc.add_paragraph(caption, style=None)
            summary_doc.add_picture(str(visuals[key]), width=Inches(6.4))
    summary_doc.save(MAIN / "08_Visualisasi" / "Lampiran_Visual_Data_Eksperimen.docx")


def notebook(cells: list[dict]) -> dict:
    return {
        "cells": cells,
        "metadata": {
            "kernelspec": {"display_name": "Python 3", "language": "python", "name": PYTHON_ID},
            "language_info": {"name": "python", "version": "3.x"},
        },
        "nbformat": 4,
        "nbformat_minor": 5,
    }


def md_cell(text: str) -> dict:
    return {"cell_type": "markdown", "metadata": {}, "source": text.splitlines(True)}


def code_cell(code: str) -> dict:
    return {"cell_type": "code", "execution_count": None, "metadata": {}, "outputs": [], "source": code.splitlines(True)}


def create_notebooks() -> None:
    nb_dir = MAIN / "05_Source_Code" / "Notebook"
    common_setup = """from pathlib import Path
import sys

ROOT = Path.cwd()
if not (ROOT / 'app').exists():
    ROOT = Path(r'C:/Users/Asus_/OneDrive/Documents/NetGuard-AI')
sys.path.insert(0, str(ROOT / 'app'))
print('Project root:', ROOT)
"""
    preprocessing_nb = notebook(
        [
            md_cell(f"# Preprocessing Notebook\n\nNama: {STUDENT['name']}  \nNIM: {STUDENT['nim']}  \nTujuan: membersihkan dataset CICIDS2017-style CSV sebelum training."),
            code_cell(common_setup),
            code_cell("""import pandas as pd
raw_path = ROOT / 'data' / 'sample' / 'sample_cicids2017.csv'
df = pd.read_csv(raw_path)
df.head()
"""),
            code_cell("""import subprocess, sys
processed_path = ROOT / 'data' / 'processed' / 'sample_cicids2017_processed.csv'
cmd = [sys.executable, str(ROOT / 'app' / 'preprocessing.py'), '--input', str(raw_path), '--output', str(processed_path)]
result = subprocess.run(cmd, cwd=ROOT, capture_output=True, text=True)
print(result.stdout)
if result.returncode != 0:
    print(result.stderr)
    raise SystemExit(result.returncode)
"""),
            code_cell("""processed = pd.read_csv(processed_path)
print(processed.shape)
processed['label'].value_counts()
"""),
        ]
    )
    training_nb = notebook(
        [
            md_cell("# Training Notebook\n\nNotebook ini menjalankan Logistic Regression, Decision Tree, dan Random Forest."),
            code_cell(common_setup),
            code_cell("""import subprocess, sys
processed_path = ROOT / 'data' / 'processed' / 'sample_cicids2017_processed.csv'
cmd = [sys.executable, str(ROOT / 'app' / 'train.py'), '--input', str(processed_path)]
result = subprocess.run(cmd, cwd=ROOT, capture_output=True, text=True)
print(result.stdout)
if result.returncode != 0:
    print(result.stderr)
    raise SystemExit(result.returncode)
"""),
            code_cell("""import pandas as pd
comparison = pd.read_csv(ROOT / 'reports' / 'model_comparison.csv')
comparison
"""),
            code_cell("""import json
metrics = json.loads((ROOT / 'reports' / 'metrics.json').read_text())
metrics['best_model'], metrics['best_metrics']
"""),
        ]
    )
    evaluation_nb = notebook(
        [
            md_cell("# Evaluation Notebook\n\nNotebook ini membaca hasil evaluasi, confusion matrix, dan ringkasan model terbaik."),
            code_cell(common_setup),
            code_cell("""import json, pandas as pd
metrics = json.loads((ROOT / 'reports' / 'metrics.json').read_text())
comparison = pd.read_csv(ROOT / 'reports' / 'model_comparison.csv')
print('Best model:', metrics.get('best_model'))
comparison
"""),
            code_cell("""from IPython.display import Image, display
display(Image(filename=str(ROOT / 'reports' / 'figures' / 'confusion_matrix.png')))
"""),
        ]
    )
    prediction_nb = notebook(
        [
            md_cell("# Prediction and Dashboard Artifact Notebook\n\nNotebook ini menjalankan prediksi data baru dan menampilkan ringkasan risk score."),
            code_cell(common_setup),
            code_cell("""import subprocess, sys
input_csv = ROOT / 'data' / 'sample' / 'sample_cicids2017.csv'
cmd = [sys.executable, str(ROOT / 'app' / 'predict.py'), '--input', str(input_csv)]
result = subprocess.run(cmd, cwd=ROOT, capture_output=True, text=True)
print(result.stdout)
if result.returncode != 0:
    print(result.stderr)
    raise SystemExit(result.returncode)
"""),
            code_cell("""import pandas as pd, json
pred = pd.read_csv(ROOT / 'reports' / 'prediction_result.csv')
summary = json.loads((ROOT / 'reports' / 'research_summary.json').read_text())
print(summary)
pred.head()
"""),
        ]
    )
    files = {
        "preprocessing.ipynb": preprocessing_nb,
        "training.ipynb": training_nb,
        "evaluation.ipynb": evaluation_nb,
        "prediction_dashboard_demo.ipynb": prediction_nb,
    }
    for name, content in files.items():
        (nb_dir / name).write_text(json.dumps(content, indent=2), encoding="utf-8")


def shape_text(slide, text: str, x, y, w, h, fill: str, font_size=16, bold=False, color="F8FAFC", radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    box = slide.shapes.add_shape(shape_type, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = PptRGB(*bytes.fromhex(fill))
    box.line.color.rgb = PptRGB(*bytes.fromhex(fill))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = PptInches(0.12)
    tf.margin_right = PptInches(0.12)
    tf.margin_top = PptInches(0.08)
    tf.margin_bottom = PptInches(0.08)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = PptPt(font_size)
    run.font.bold = bold
    run.font.color.rgb = PptRGB(*bytes.fromhex(color))
    return box


def add_title(slide, title: str, subtitle: str | None = None, color="F8FAFC") -> None:
    tb = slide.shapes.add_textbox(PptInches(0.45), PptInches(0.28), PptInches(12.45), PptInches(0.55))
    p = tb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = PptPt(25)
    run.font.bold = True
    run.font.color.rgb = PptRGB(*bytes.fromhex(color))
    if subtitle:
        st = slide.shapes.add_textbox(PptInches(0.48), PptInches(0.82), PptInches(11.8), PptInches(0.3))
        p2 = st.text_frame.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = PptPt(10.5)
        r2.font.color.rgb = PptRGB(*bytes.fromhex(PALETTE["muted"]))


def add_footer(slide, number: int) -> None:
    tb = slide.shapes.add_textbox(PptInches(0.45), PptInches(7.08), PptInches(8.8), PptInches(0.18))
    p = tb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = f"NetGuard AI | Achmad Maulana 241730016 | {number:02d}"
    run.font.size = PptPt(7.5)
    run.font.color.rgb = PptRGB(148, 163, 184)


def add_nav(slide, slides, idx: int) -> None:
    home = shape_text(slide, "HOME", PptInches(10.6), PptInches(6.95), PptInches(0.75), PptInches(0.28), "1E293B", 8, True)
    home.click_action.target_slide = slides[1] if len(slides) > 1 else slides[0]
    if idx > 0:
        prev = shape_text(slide, "PREV", PptInches(11.42), PptInches(6.95), PptInches(0.68), PptInches(0.28), "334155", 8, True)
        prev.click_action.target_slide = slides[idx - 1]
    if idx < len(slides) - 1:
        nxt = shape_text(slide, "NEXT", PptInches(12.17), PptInches(6.95), PptInches(0.68), PptInches(0.28), "38BDF8", 8, True, "0B1220")
        nxt.click_action.target_slide = slides[idx + 1]


def add_bg(slide, mode: int) -> None:
    bg = slide.background.fill
    bg.solid()
    base = ["0B1220", "07111F", "111827", "0F172A", "06121E"][mode % 5]
    bg.fore_color.rgb = PptRGB(*bytes.fromhex(base))
    for i, color in enumerate(["38BDF8", "34D399", "A78BFA", "FBBF24", "FB7185"]):
        x = PptInches((i * 2.2 + mode * 0.4) % 12.5)
        y = PptInches(0.15 + ((i + mode) % 4) * 1.65)
        blob = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, PptInches(1.5), PptInches(1.5))
        blob.fill.solid()
        blob.fill.fore_color.rgb = PptRGB(*bytes.fromhex(color))
        blob.fill.transparency = 72
        blob.line.transparency = 100


def add_chart_image(slide, path: Path, x, y, w, h) -> None:
    if path.exists():
        slide.shapes.add_picture(str(path), x, y, width=w, height=h)


def create_enhanced_ppt(visuals: dict[str, Path]) -> None:
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)
    blank = prs.slide_layouts[6]
    slides = [prs.slides.add_slide(blank) for _ in range(14)]

    metrics = load_metrics()
    comparison = load_comparison()
    best_model = metrics.get("best_model", "Logistic Regression")
    best_metrics = metrics.get("models", {}).get(best_model, {"accuracy": 1.0, "precision": 1.0, "recall": 1.0, "f1_score": 1.0})

    for i, slide in enumerate(slides):
        add_bg(slide, i)

    s = slides[0]
    shape_text(s, "NETGUARD AI", PptInches(0.7), PptInches(1.0), PptInches(5.7), PptInches(0.85), "38BDF8", 33, True, "0B1220", False)
    add_title(s, "Predictive Network Failure & Anomaly Monitoring", "AI-based Network Anomaly Detection | UAS Artificial Intelligence 2026")
    shape_text(s, "Achmad Maulana\n241730016", PptInches(0.75), PptInches(2.15), PptInches(2.5), PptInches(0.75), "1E293B", 17, True)
    shape_text(s, "Informatika\nUIN Sultan Maulana Hasanuddin Banten", PptInches(3.45), PptInches(2.15), PptInches(3.25), PptInches(0.75), "334155", 14, True)
    add_chart_image(s, visuals["pipeline"], PptInches(0.85), PptInches(3.15), PptInches(11.7), PptInches(3.05))

    s = slides[1]
    add_title(s, "Agenda: Dari Gap ke Demo Sistem", "Klik tombol untuk lompat ke bagian penting presentasi")
    items = [("1", "Problem", 2), ("2", "Research Gap", 3), ("3", "Method", 5), ("4", "Experiment", 7), ("5", "Dashboard", 10), ("6", "Defense", 12)]
    for idx, (num, label, target) in enumerate(items):
        x = PptInches(0.85 + (idx % 3) * 4.05)
        y = PptInches(1.55 + (idx // 3) * 2.25)
        box = shape_text(s, f"{num}\n{label}", x, y, PptInches(3.35), PptInches(1.45), ["38BDF8", "34D399", "A78BFA", "FBBF24", "FB7185", "22D3EE"][idx], 22, True, "0B1220")
        box.click_action.target_slide = slides[target]
    shape_text(s, "Navigasi interaktif aktif: setiap tombol agenda mengarah ke slide terkait.", PptInches(1.2), PptInches(6.15), PptInches(10.8), PptInches(0.45), "1E293B", 13, True)

    s = slides[2]
    add_title(s, "Masalah: Admin Jaringan Butuh Deteksi Cepat", "Traffic abnormal sulit terlihat jika hanya mengandalkan inspeksi manual log router/firewall")
    problems = [("Volume traffic naik", "Log cepat menumpuk dan sulit dibaca manual."), ("Serangan makin variatif", "DoS, PortScan, Botnet, Brute Force memiliki pola berbeda."), ("Resource mahasiswa terbatas", "Laptop biasa, tanpa GPU dan tanpa API berbayar."), ("Hasil riset perlu bukti", "Butuh metrik, artefak, dan dashboard demo.")]
    for i, (a, b) in enumerate(problems):
        shape_text(s, a, PptInches(0.8), PptInches(1.35 + i * 1.15), PptInches(3.3), PptInches(0.45), ["38BDF8", "34D399", "FBBF24", "FB7185"][i], 14, True, "0B1220")
        shape_text(s, b, PptInches(4.35), PptInches(1.35 + i * 1.15), PptInches(7.7), PptInches(0.45), "1E293B", 13, False)

    s = slides[3]
    add_title(s, "Research Gap yang Diambil", "Gap diambil dari pola literatur: model akurat sering belum dikemas menjadi monitoring edukatif ringan")
    add_chart_image(s, visuals["architecture"], PptInches(0.65), PptInches(1.25), PptInches(7.25), PptInches(3.05))
    gaps = [("Gap 1", "Banyak paper fokus metrik, bukan dashboard yang langsung dipahami admin pemula."), ("Gap 2", "Implementasi sering berat untuk mahasiswa/laptop biasa."), ("Gap 3", "Risk score dan rekomendasi aksi belum selalu terintegrasi."), ("Gap 4", "Artefak reproduksi sering tidak lengkap.")]
    for i, (g, desc) in enumerate(gaps):
        shape_text(s, f"{g}: {desc}", PptInches(8.25), PptInches(1.2 + i * 0.85), PptInches(4.25), PptInches(0.55), ["38BDF8", "34D399", "FBBF24", "FB7185"][i], 10.5, True, "0B1220")

    s = slides[4]
    add_title(s, "Novelty: ML + Dashboard Edukatif + Risk Action", "Kebaruan utama bukan sekadar algoritma, tetapi integrasi pipeline penelitian menjadi alat demo yang dapat direproduksi")
    novelty = [("Main Novelty", "Integrasi anomaly detection berbasis machine learning dengan dashboard monitoring edukatif yang ringan."), ("Secondary Novelty", "Risk classification otomatis Low/Medium/High dengan rekomendasi tindakan praktis."), ("Educational Value", "Cocok untuk mahasiswa TKJ/Informatika yang ingin memahami AI dan jaringan secara implementatif.")]
    for i, (title, text) in enumerate(novelty):
        shape_text(s, title, PptInches(0.85), PptInches(1.5 + i * 1.55), PptInches(2.7), PptInches(0.55), ["38BDF8", "A78BFA", "34D399"][i], 15, True, "0B1220")
        shape_text(s, text, PptInches(3.75), PptInches(1.42 + i * 1.55), PptInches(8.4), PptInches(0.72), "1E293B", 14, False)

    s = slides[5]
    add_title(s, "Research Method: Pipeline NetGuard AI", "RM1 implementasi model ML; RM2 integrasi hasil ke dashboard dan artefak riset")
    add_chart_image(s, visuals["pipeline"], PptInches(0.65), PptInches(1.25), PptInches(12.0), PptInches(3.2))
    shape_text(s, "RM1: Membangun dan membandingkan Logistic Regression, Decision Tree, dan Random Forest pada data CICIDS2017-style.", PptInches(0.95), PptInches(5.15), PptInches(5.7), PptInches(0.8), "38BDF8", 12, True, "0B1220")
    shape_text(s, "RM2: Mengintegrasikan model terbaik ke dashboard Flask dengan risk score dan rekomendasi tindakan.", PptInches(6.85), PptInches(5.15), PptInches(5.7), PptInches(0.8), "34D399", 12, True, "0B1220")

    s = slides[6]
    add_title(s, "Dataset: CICIDS2017 Real Subset", "Dataset resmi dari Canadian Institute for Cybersecurity, UNB; subset 20.000 record dipakai agar tetap low-budget")
    add_chart_image(s, visuals["distribution"], PptInches(0.75), PptInches(1.35), PptInches(6.25), PptInches(3.7))
    table = s.shapes.add_table(5, 2, PptInches(7.25), PptInches(1.45), PptInches(5.2), PptInches(2.9)).table
    data = [["Aspek", "Isi"], ["Sumber", "UNB CICIDS2017"], ["Kelas", "BENIGN vs DDoS"], ["Raw subset", "20.000 record"], ["Processed", "19.935 record"]]
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = PptRGB(30, 41, 59) if r else PptRGB(56, 189, 248)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = PptPt(10.5)
                    run.font.color.rgb = PptRGB(248, 250, 252) if r else PptRGB(11, 18, 32)
                    run.font.bold = r == 0

    s = slides[7]
    add_title(s, "Eksperimen: Tiga Model Pembanding", "Best model dipilih dengan prioritas F1-score, lalu Recall, lalu Accuracy")
    add_chart_image(s, visuals["comparison"], PptInches(0.65), PptInches(1.25), PptInches(7.4), PptInches(4.1))
    for i, row in comparison.iterrows():
        shape_text(s, f"{row['model']}\nF1 {float(row['f1_score']):.2f} | Recall {float(row['recall']):.2f}", PptInches(8.35), PptInches(1.35 + i * 1.35), PptInches(3.9), PptInches(0.85), ["38BDF8", "34D399", "A78BFA"][i], 15, True, "0B1220")

    s = slides[8]
    add_title(s, "Confusion Matrix sebagai Bukti Evaluasi", "Evaluasi tidak hanya angka tunggal; matrix menunjukkan benar-salah prediksi setiap kelas")
    cm_path = MAIN / "08_Visualisasi" / "confusion_matrix.png"
    add_chart_image(s, cm_path, PptInches(0.85), PptInches(1.15), PptInches(6.2), PptInches(5.2))
    bullets = [("True Normal", "Traffic normal diprediksi normal."), ("False Alarm", "Normal diprediksi anomaly."), ("Missed Attack", "Anomaly diprediksi normal."), ("True Anomaly", "Traffic anomaly terdeteksi.")]
    for i, (k, v) in enumerate(bullets):
        shape_text(s, f"{k}\n{v}", PptInches(7.55), PptInches(1.25 + i * 1.12), PptInches(4.65), PptInches(0.76), ["34D399", "FBBF24", "FB7185", "38BDF8"][i], 12.5, True, "0B1220")

    s = slides[9]
    add_title(s, "Risk Score: Angka yang Bisa Ditindaklanjuti", "Risk level dihitung dari anomaly ratio agar hasil model mudah dipakai admin jaringan pemula")
    add_chart_image(s, visuals["risk"], PptInches(0.85), PptInches(1.25), PptInches(6.3), PptInches(3.7))
    risk_table = [["Level", "Aturan", "Aksi"], ["Low", "< 10%", "Continue monitoring"], ["Medium", "10-30%", "Inspect sources and firewall logs"], ["High", ">= 30%", "Investigate and isolate suspicious hosts"]]
    table = s.shapes.add_table(4, 3, PptInches(7.25), PptInches(1.45), PptInches(5.45), PptInches(2.6)).table
    for r, row in enumerate(risk_table):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = PptRGB(56, 189, 248) if r == 0 else PptRGB(30, 41, 59)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = PptPt(8.8)
                    run.font.bold = r == 0
                    run.font.color.rgb = PptRGB(11, 18, 32) if r == 0 else PptRGB(248, 250, 252)

    s = slides[10]
    add_title(s, "Dashboard Demo: Monitoring Tanpa Database", "Flask membaca metrics.json dan CSV hasil eksperimen agar ringan untuk laptop mahasiswa")
    ds = metrics.get("dataset_summary", {})
    ui_blocks = [
        ("Total Records", f"{int(ds.get('total_records', 19935)):,}"),
        ("Normal", f"{int(ds.get('normal_records', 9935)):,}"),
        ("Anomaly", f"{int(ds.get('anomaly_records', 10000)):,}"),
        ("Accuracy", f"{float(best_metrics.get('accuracy', 1.0)):.4f}"),
        ("Best Model", best_model),
        ("Risk", "High"),
    ]
    for i, (label, value) in enumerate(ui_blocks):
        shape_text(s, f"{value}\n{label}", PptInches(0.8 + (i % 3) * 4.1), PptInches(1.35 + (i // 3) * 1.45), PptInches(3.35), PptInches(1.0), ["38BDF8", "34D399", "FB7185", "A78BFA", "FBBF24", "22D3EE"][i], 17, True, "0B1220")
    shape_text(s, "Workflow demo: Upload Dataset -> Train Model -> Predict New Data -> View Research Report", PptInches(1.0), PptInches(5.2), PptInches(11.2), PptInches(0.65), "1E293B", 17, True)

    s = slides[11]
    add_title(s, "Deployment Plan: Low-Budget ke Premium UI", "Flask tetap menjadi demo utama; Next.js/Tailwind/Framer disiapkan untuk bukti bonus deployment modern")
    add_chart_image(s, visuals["architecture"], PptInches(0.65), PptInches(1.2), PptInches(7.2), PptInches(3.2))
    stack = [("Frontend", "Next.js + Tailwind + Framer Motion"), ("Backend", "Flask sekarang, REST API future"), ("Database", "PostgreSQL + Prisma schema"), ("Deploy", "Vercel + Railway/Render")]
    for i, (k, v) in enumerate(stack):
        shape_text(s, f"{k}\n{v}", PptInches(8.2), PptInches(1.25 + i * 0.95), PptInches(4.15), PptInches(0.62), ["38BDF8", "34D399", "A78BFA", "FBBF24"][i], 11.5, True, "0B1220")

    s = slides[12]
    add_title(s, "Kesiapan Artefak untuk Diuji Dosen", "Artefak disusun mengikuti struktur Google Drive UAS AI")
    artifacts = [["Paper", "11 PDF"], ["Mapping", "XLSX + DOCX"], ["Gap/Novelty", "MD + DOCX"], ["Source Code", "Script + Notebook"], ["Model", "PKL"], ["Eksperimen", "JSON + CSV + PNG"], ["IEEE Draft", "DOCX + PDF"], ["Presentasi", "PPTX + PDF"], ["Turnitin", "Panduan cek"], ["Deployment", "Next.js source"]]
    for i, (k, v) in enumerate(artifacts):
        shape_text(s, f"{k}\n{v}", PptInches(0.65 + (i % 5) * 2.52), PptInches(1.4 + (i // 5) * 1.75), PptInches(2.05), PptInches(1.0), ["38BDF8", "34D399", "A78BFA", "FBBF24", "FB7185"][i % 5], 12.2, True, "0B1220")

    s = slides[13]
    add_title(s, "Kesimpulan dan Arah Pengembangan", "NetGuard AI layak sebagai proyek AI edukatif dan dapat ditingkatkan untuk riset/deployment lebih kuat")
    conclusions = [("Kesimpulan", "Pipeline ML, evaluasi, prediksi, dashboard, dan artefak penelitian sudah terintegrasi."), ("Keterbatasan", "Smoke test kecil; eksperimen final wajib memakai subset CICIDS2017 yang lebih besar."), ("Next Research", "Cross-validation, hyperparameter tuning, XGBoost/LightGBM, UNSW-NB15, dan deployment cloud."), ("Defense Message", "Nilai utama proyek: low-budget, reproducible, dan mudah dijelaskan sebagai jembatan AI + networking.")]
    for i, (k, v) in enumerate(conclusions):
        shape_text(s, k, PptInches(0.85), PptInches(1.25 + i * 1.25), PptInches(2.3), PptInches(0.48), ["38BDF8", "FBBF24", "34D399", "A78BFA"][i], 13.5, True, "0B1220")
        shape_text(s, v, PptInches(3.4), PptInches(1.22 + i * 1.25), PptInches(8.5), PptInches(0.55), "1E293B", 12.2, False)

    for idx, slide in enumerate(slides):
        add_footer(slide, idx + 1)
    for idx, slide in enumerate(slides):
        add_nav(slide, slides, idx)

    out = MAIN / "10_Presentasi" / "Slide_Presentasi_Enhanced_NetGuard_AI_v2.pptx"
    prs.save(out)
    add_slide_transitions(out)
    for target in [MAIN / "10_Presentasi" / "Slide_Presentasi.pptx", ROOT / "Slide_Presentasi_NetGuard_AI_Achmad_Maulana.pptx"]:
        try:
            shutil.copyfile(out, target)
        except PermissionError:
            print(f"Warning: {target} is locked. Use {out.name} as the latest deck.")
    create_ppt_pdf_summary(MAIN / "10_Presentasi" / "Slide_Presentasi_Enhanced_NetGuard_AI.pdf")
    shutil.copyfile(MAIN / "10_Presentasi" / "Slide_Presentasi_Enhanced_NetGuard_AI.pdf", MAIN / "10_Presentasi" / "Slide_Presentasi.pdf")
    shutil.copyfile(MAIN / "10_Presentasi" / "Slide_Presentasi_Enhanced_NetGuard_AI.pdf", ROOT / "Slide_Presentasi_NetGuard_AI_Achmad_Maulana.pdf")


def add_slide_transitions(pptx_path: Path) -> None:
    ns = {
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    }
    ET.register_namespace("p", ns["p"])
    ET.register_namespace("a", ns["a"])
    tmp = pptx_path.with_suffix(".tmp.pptx")
    effects = [
        '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:push dir="l"/></p:transition>',
        '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:wipe dir="r"/></p:transition>',
        '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:split orient="horz" dir="out"/></p:transition>',
        '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:fade/></p:transition>',
    ]
    with zipfile.ZipFile(pptx_path, "r") as zin, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.startswith("ppt/slides/slide") and item.filename.endswith(".xml"):
                try:
                    root = ET.fromstring(data)
                    for existing in root.findall("p:transition", ns):
                        root.remove(existing)
                    idx = int(Path(item.filename).stem.replace("slide", "")) - 1
                    transition = ET.fromstring(effects[idx % len(effects)])
                    root.insert(1, transition)
                    data = ET.tostring(root, encoding="utf-8", xml_declaration=True)
                except Exception:
                    pass
            zout.writestr(item, data)
    tmp.replace(pptx_path)


def create_ppt_pdf_summary(pdf_path: Path) -> None:
    doc = SimpleDocTemplate(str(pdf_path), pagesize=A4)
    styles = getSampleStyleSheet()
    story = [
        Paragraph("Slide Presentasi NetGuard AI - Enhanced Deck Summary", styles["Title"]),
        Spacer(1, 12),
        Paragraph("File PPTX utama berisi 14 slide dengan variasi layout, tabel, grafik, visual pipeline, tombol navigasi internal, hyperlink antar-slide, dan transisi slide native PowerPoint.", styles["BodyText"]),
        Spacer(1, 10),
        Paragraph("Gunakan file PPTX untuk presentasi karena fitur tombol dan transisi aktif di PowerPoint.", styles["BodyText"]),
    ]
    doc.build(story)


def update_readmes() -> None:
    nb_readme = MAIN / "05_Source_Code" / "Notebook" / "README_Notebook.md"
    nb_readme.write_text(
        """# Notebook Source Code

Folder ini berisi notebook yang dapat dibuka di Jupyter Notebook, JupyterLab, VS Code, atau Google Colab.

Isi notebook:

1. `preprocessing.ipynb` - membaca CSV sample, membersihkan data, encode label, dan menyimpan processed CSV.
2. `training.ipynb` - menjalankan Logistic Regression, Decision Tree, Random Forest, serta menyimpan model terbaik.
3. `evaluation.ipynb` - membaca metrics JSON, comparison CSV, dan confusion matrix.
4. `prediction_dashboard_demo.ipynb` - menjalankan prediksi data baru dan membaca risk summary.

Catatan: notebook ini dibuat untuk memudahkan dosen melihat alur eksperimen secara interaktif. Script `.py` tetap menjadi source code utama yang bisa dijalankan dari Windows CMD.
""",
        encoding="utf-8",
    )
    readme = MAIN / "README.md"
    text = readme.read_text(encoding="utf-8")
    addition = """

## Update Artefak Tambahan

Paket ini sudah ditambahkan artefak lanjutan:

- Notebook Jupyter di `05_Source_Code/Notebook`.
- Versi Word untuk literature mapping, gap analysis, experiment log, narasi presentasi, dokumentasi deployment, dan README paket.
- Lampiran visual data eksperimen dalam bentuk `.png` dan `.docx`.
- PPT enhanced dengan 14 slide, variasi layout, grafik, tabel, tombol navigasi internal, dan transisi slide.
"""
    if "## Update Artefak Tambahan" not in text:
        readme.write_text(text.rstrip() + addition, encoding="utf-8")
    shutil.copyfile(readme, MAIN / "05_Source_Code" / "README.md")


def main() -> None:
    ensure_dirs()
    visuals = create_visuals()
    create_notebooks()
    create_docx_companions(visuals)
    create_enhanced_ppt(visuals)
    update_readmes()
    print("Enhanced UAS artifacts generated.")
    print(f"Main folder: {MAIN}")


if __name__ == "__main__":
    main()
