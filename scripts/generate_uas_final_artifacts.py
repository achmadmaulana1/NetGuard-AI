from pathlib import Path
import json
import shutil
from datetime import date

from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT, WD_CELL_VERTICAL_ALIGNMENT
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt
from pptx.dml.color import RGBColor as PptRGB
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY, TA_LEFT
from reportlab.lib.pagesizes import letter, landscape
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import (
    SimpleDocTemplate,
    Paragraph,
    Spacer,
    Table,
    TableStyle,
    PageBreak,
)


ROOT = Path(__file__).resolve().parents[1]
STUDENT = "Achmad Maulana"
NIM = "241730016"
PROGRAM = "Program Studi Informatika"
FACULTY = "Fakultas Sains dan Teknologi"
UNIVERSITY = "Universitas Islam Negeri Sultan Maulana Hasanuddin Banten"
YEAR = "2026"
MAIN = ROOT / "Achmad_Maulana_241730016_UAS_AI"
REPORTS = ROOT / "reports"


PAPERS = [
    {
        "no": 1,
        "year": 2021,
        "authors": "Z. K. Maseer, R. Yusof, N. Bahaman, S. A. Mostafa, and C. F. M. Foozy",
        "title": "Benchmarking of Machine Learning for Anomaly Based Intrusion Detection Systems in the CICIDS2017 Dataset",
        "publisher": "IEEE Access",
        "category": "IEEE / Scopus-indexed",
        "doi": "10.1109/ACCESS.2021.3056614",
        "link": "https://doi.org/10.1109/ACCESS.2021.3056614",
        "dataset": "CICIDS2017",
        "method": "ML benchmarking",
        "result": "Multiple ML models benchmarked for anomaly-based IDS.",
        "acc": "Reported per model/attack in paper",
        "precision": "Reported in paper",
        "recall": "Reported in paper",
        "f1": "Reported in paper",
        "strength": "Strong baseline reference for classical ML on CICIDS2017.",
        "weakness": "Focuses on model benchmarking, not educational dashboard deployment.",
        "gap": "Need practical dashboard and reproducible student-ready implementation.",
    },
    {
        "no": 2,
        "year": 2024,
        "authors": "R. Dube",
        "title": "Faulty Use of the CIC-IDS 2017 Dataset in Information Security Research",
        "publisher": "Journal of Computer Virology and Hacking Techniques",
        "category": "Springer / Scopus-indexed placeholder",
        "doi": "10.1007/s11416-023-00509-7",
        "link": "https://link.springer.com/article/10.1007/s11416-023-00509-7",
        "dataset": "CICIDS2017",
        "method": "Dataset validity analysis",
        "result": "Highlights methodological risks when CICIDS2017 is used incorrectly.",
        "acc": "-",
        "precision": "-",
        "recall": "-",
        "f1": "-",
        "strength": "Important validity warning for CICIDS2017 research.",
        "weakness": "Does not provide a deployed anomaly monitoring prototype.",
        "gap": "Need transparent preprocessing and clear reproducibility documentation.",
    },
    {
        "no": 3,
        "year": 2024,
        "authors": "M. Sajid, K. R. Malik, A. Almogren, T. S. Malik, A. H. Khan, J. Tanveer, and A. U. Rehman",
        "title": "Enhancing Intrusion Detection: A Hybrid Machine and Deep Learning Approach",
        "publisher": "Journal of Cloud Computing",
        "category": "SpringerOpen / Scopus-indexed placeholder",
        "doi": "10.1186/s13677-024-00685-x",
        "link": "https://link.springer.com/article/10.1186/s13677-024-00685-x",
        "dataset": "CICIDS2017, UNSW-NB15, NSL-KDD, WSN-DS",
        "method": "XGBoost, CNN, LSTM, CNN-LSTM",
        "result": "Hybrid ML/DL models evaluated on multiple IDS datasets.",
        "acc": "Approx. 97.90%-98.40% for CICIDS2017 binary tasks as reported by source",
        "precision": "See paper",
        "recall": "See paper",
        "f1": "See paper",
        "strength": "Strong multi-dataset and hybrid-model evaluation.",
        "weakness": "More complex than a low-budget student deployment.",
        "gap": "Need lightweight implementation that can run on normal laptops.",
    },
    {
        "no": 4,
        "year": 2025,
        "authors": "M. T. Abdelaziz, A. Radwan, H. Mamdouh, A. S. Saad, A. S. Abuzaid, A. A. AbdElhakeem, S. Zakzouk, K. Moussa, and M. S. Darweesh",
        "title": "Enhancing Network Threat Detection with Random Forest-Based NIDS and Permutation Feature Importance",
        "publisher": "Journal of Network and Systems Management",
        "category": "Springer / Scopus-indexed placeholder",
        "doi": "10.1007/s10922-024-09874-0",
        "link": "https://link.springer.com/article/10.1007/s10922-024-09874-0",
        "dataset": "CICIDS2017",
        "method": "Random Forest, permutation feature importance",
        "result": "High F1 with feature importance analysis.",
        "acc": "See paper",
        "precision": "See paper",
        "recall": "See paper",
        "f1": "Weighted F1 and macro F1 reported by paper",
        "strength": "Adds interpretability through feature importance.",
        "weakness": "Deployment and beginner workflow are not central contributions.",
        "gap": "Need dashboard-based artifact and operational recommendations.",
    },
    {
        "no": 5,
        "year": 2024,
        "authors": "M. A. Talukder, M. M. Islam, M. A. Uddin, K. F. Hasan, S. Sharmin, S. A. Alyami, and M. A. Moni",
        "title": "Machine Learning-Based Network Intrusion Detection for Big and Imbalanced Data Using Oversampling, Stacking Feature Embedding and Feature Extraction",
        "publisher": "Journal of Big Data",
        "category": "SpringerOpen / Scopus-indexed placeholder",
        "doi": "10.1186/s40537-024-00886-w",
        "link": "https://link.springer.com/article/10.1186/s40537-024-00886-w",
        "dataset": "UNSW-NB15, CICIDS2017, CICIDS2018",
        "method": "Oversampling, stacking feature embedding, PCA, DT/RF/ET",
        "result": "Tree-based models reach very high accuracy on benchmark datasets.",
        "acc": "CICIDS2017 up to 99.99% for some tree-based models as reported by source",
        "precision": "See paper",
        "recall": "See paper",
        "f1": "See paper",
        "strength": "Strong handling of large and imbalanced IDS data.",
        "weakness": "Pipeline is more complex and dashboard is not the emphasis.",
        "gap": "Need simpler low-budget framework with clear artifact packaging.",
    },
    {
        "no": 6,
        "year": 2024,
        "authors": "M. Deng, C. Sun, Y. Kan, and S. Fan",
        "title": "Network Intrusion Detection Based on Deep Belief Network Broad Equalization Learning System",
        "publisher": "Electronics",
        "category": "MDPI / Scopus-indexed placeholder",
        "doi": "10.3390/electronics13153014",
        "link": "https://www.mdpi.com/2079-9292/13/15/3014",
        "dataset": "CICIDS2017 and related IDS datasets",
        "method": "DBN and broad equalization learning",
        "result": "Deep feature learning for high-dimensional intrusion data.",
        "acc": "See paper",
        "precision": "See paper",
        "recall": "See paper",
        "f1": "See paper",
        "strength": "Handles high-dimensional data with deep representation.",
        "weakness": "Higher complexity and lower beginner explainability.",
        "gap": "Need explainable and teachable implementation for undergraduate context.",
    },
    {
        "no": 7,
        "year": 2024,
        "authors": "R. Mohammad, F. Saeed, A. A. Almazroi, F. S. Alsubaei, and A. A. Almazroi",
        "title": "Enhancing Intrusion Detection Systems Using a Deep Learning and Data Augmentation Approach",
        "publisher": "Systems",
        "category": "MDPI / Scopus-indexed placeholder",
        "doi": "10.3390/systems12030079",
        "link": "https://www.mdpi.com/2079-8954/12/3/79",
        "dataset": "Augmented CIC-IDS-2017",
        "method": "Deep learning and data augmentation",
        "result": "Reported accuracy up to about 91% on augmented CIC-IDS-2017.",
        "acc": "Up to about 91% as reported by source page",
        "precision": "See paper",
        "recall": "See paper",
        "f1": "See paper",
        "strength": "Uses augmentation to improve IDS modeling.",
        "weakness": "Does not target lightweight dashboard deployment.",
        "gap": "Need UI-integrated and reproducible student research artifact.",
    },
    {
        "no": 8,
        "year": 2024,
        "authors": "M. Wang, N. Yang, Y. Guo, and N. Weng",
        "title": "Learn-IDS: Bridging Gaps between Datasets and Learning-Based Network Intrusion Detection",
        "publisher": "Electronics",
        "category": "MDPI / Scopus-indexed placeholder",
        "doi": "10.3390/electronics13061072",
        "link": "https://www.mdpi.com/2079-9292/13/6/1072",
        "dataset": "Multiple IDS datasets",
        "method": "Learning-based IDS framework",
        "result": "Framework to bridge dataset and learning-process gaps.",
        "acc": "See paper",
        "precision": "See paper",
        "recall": "See paper",
        "f1": "See paper",
        "strength": "Directly discusses dataset-learning gaps.",
        "weakness": "Less focused on low-budget educational deployment.",
        "gap": "Need practical student-ready system packaging.",
    },
    {
        "no": 9,
        "year": 2025,
        "authors": "V. Pai, K. Pai, S. Manjunatha, S. Hirmeti, and V. V. Bhat",
        "title": "Adaptive Network Anomaly Detection Using Machine Learning Approaches",
        "publisher": "EURASIP Journal on Information Security",
        "category": "SpringerOpen / Scopus-indexed placeholder",
        "doi": "10.1186/s13635-025-00216-4",
        "link": "https://link.springer.com/article/10.1186/s13635-025-00216-4",
        "dataset": "Network anomaly datasets",
        "method": "KNN, ensemble and ML approaches",
        "result": "Adaptive machine-learning approach for anomaly detection.",
        "acc": "See paper",
        "precision": "See paper",
        "recall": "See paper",
        "f1": "See paper",
        "strength": "Adaptive framing and recent source.",
        "weakness": "Need to verify exact dataset and metrics for comparison.",
        "gap": "Need implementation aligned with undergraduate UAS constraints.",
    },
    {
        "no": 10,
        "year": 2025,
        "authors": "U. Ahmed, M. Nazir, A. Sarwar, T. Ali, E.-H. M. Aggoune, T. Shahzad, and M. A. Khan",
        "title": "Signature-Based Intrusion Detection Using Machine Learning and Deep Learning Approaches Empowered with Fuzzy Clustering",
        "publisher": "Scientific Reports",
        "category": "Nature Portfolio / Scopus-indexed placeholder",
        "doi": "10.1038/s41598-025-85866-7",
        "link": "https://www.nature.com/articles/s41598-025-85866-7",
        "dataset": "IDS datasets including CICIDS2017 context",
        "method": "ML, DL, fuzzy clustering",
        "result": "Combines signature-based IDS with ML/DL and fuzzy clustering.",
        "acc": "See paper",
        "precision": "See paper",
        "recall": "See paper",
        "f1": "See paper",
        "strength": "Recent high-visibility journal source.",
        "weakness": "More complex than lightweight educational dashboard.",
        "gap": "Need simpler deployable prototype for campus learning context.",
    },
]


EXPERIMENT_ROWS = [
    ["Logistic Regression", "1.0000*", "1.0000*", "1.0000*", "1.0000*", "Baseline linear"],
    ["Decision Tree", "1.0000*", "1.0000*", "1.0000*", "1.0000*", "Baseline tree"],
    ["Random Forest", "1.0000*", "1.0000*", "1.0000*", "1.0000*", "Proposed model"],
]


def ensure_dirs():
    dirs = [
        "01_Paper",
        "02_Literature_Mapping",
        "03_Gap_Analysis",
        "04_Dataset/Raw_Dataset",
        "04_Dataset/Processed_Dataset",
        "05_Source_Code/Notebook",
        "05_Source_Code/Script",
        "06_Model",
        "07_Hasil_Eksperimen",
        "08_Visualisasi",
        "09_Draft_IEEE",
        "10_Presentasi",
        "11_Turnitin",
        "12_Deployment",
    ]
    for d in dirs:
        (MAIN / d).mkdir(parents=True, exist_ok=True)


def copy_if_exists(src, dst):
    src = Path(src)
    dst = Path(dst)
    if src.exists():
        dst.parent.mkdir(parents=True, exist_ok=True)
        shutil.copy2(src, dst)


def set_cell_shading(cell, fill):
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), fill)
    tc_pr.append(shd)


def set_doc_styles(doc):
    normal = doc.styles["Normal"]
    normal.font.name = "Times New Roman"
    normal._element.rPr.rFonts.set(qn("w:eastAsia"), "Times New Roman")
    normal.font.size = Pt(10)
    normal.font.color.rgb = RGBColor(0, 0, 0)
    for style_name, size in [("Heading 1", 14), ("Heading 2", 12), ("Heading 3", 11)]:
        style = doc.styles[style_name]
        style.font.name = "Times New Roman"
        style._element.rPr.rFonts.set(qn("w:eastAsia"), "Times New Roman")
        style.font.bold = True
        style.font.size = Pt(size)
        style.font.color.rgb = RGBColor(0, 0, 0)


def add_doc_table(doc, headers, rows, widths=None):
    table = doc.add_table(rows=1, cols=len(headers))
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.style = "Table Grid"
    hdr = table.rows[0].cells
    for i, h in enumerate(headers):
        hdr[i].text = h
        set_cell_shading(hdr[i], "D9EAF7")
        for p in hdr[i].paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.name = "Times New Roman"
                r.font.size = Pt(8)
    for row in rows:
        cells = table.add_row().cells
        for i, value in enumerate(row):
            cells[i].text = str(value)
            cells[i].vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
            for p in cells[i].paragraphs:
                for r in p.runs:
                    r.font.name = "Times New Roman"
                    r.font.size = Pt(8)
    if widths:
        for row in table.rows:
            for idx, width in enumerate(widths):
                row.cells[idx].width = Inches(width)
    doc.add_paragraph("")
    return table


def create_literature_xlsx():
    wb = Workbook()
    ws = wb.active
    ws.title = "Literature Mapping"
    headers = [
        "No",
        "Year",
        "Authors",
        "Title",
        "Publisher",
        "Category",
        "DOI",
        "Dataset",
        "Method",
        "Result",
        "Accuracy/F1",
        "Strength",
        "Weakness",
        "Research Gap",
    ]
    ws.append(headers)
    for p in PAPERS:
        ws.append([
            p["no"],
            p["year"],
            p["authors"],
            p["title"],
            p["publisher"],
            p["category"],
            p["doi"],
            p["dataset"],
            p["method"],
            p["result"],
            p["acc"] + " / " + p["f1"],
            p["strength"],
            p["weakness"],
            p["gap"],
        ])
    ws2 = wb.create_sheet("Comparison Matrix")
    ws2.append(["Article", "Dataset", "Method", "Result", "Weakness", "Research Limitation"])
    for p in PAPERS:
        ws2.append([p["title"], p["dataset"], p["method"], p["result"], p["weakness"], p["gap"]])
    ws3 = wb.create_sheet("Gap Analysis")
    ws3.append(["Gap Type", "Evidence From Literature", "Impact", "NetGuard AI Response"])
    gap_rows = [
        ["Dataset Gap", "CICIDS2017 is common but vulnerable to faulty use and single-dataset bias.", "Metrics may overstate real-world performance.", "Transparent preprocessing, explicit limitations, and future multi-dataset validation."],
        ["Method Gap", "Hybrid DL models are strong but complex.", "Low-budget students may not reproduce them.", "Classical ML baseline and Random Forest proposed model."],
        ["Deployment Gap", "Many works stop at metrics.", "Outputs are hard for beginner admins to interpret.", "Flask dashboard, risk score, and recommendations."],
        ["Evaluation Gap", "High accuracy papers do not always discuss operational meaning.", "Weak decision support.", "Metrics plus confusion matrix and risk-level interpretation."],
    ]
    for row in gap_rows:
        ws3.append(row)

    for sheet in wb.worksheets:
        for cell in sheet[1]:
            cell.font = Font(bold=True, color="FFFFFF")
            cell.fill = PatternFill("solid", fgColor="1F4E79")
            cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        for row in sheet.iter_rows():
            for cell in row:
                cell.alignment = Alignment(vertical="top", wrap_text=True)
                cell.border = Border(
                    left=Side(style="thin", color="B7B7B7"),
                    right=Side(style="thin", color="B7B7B7"),
                    top=Side(style="thin", color="B7B7B7"),
                    bottom=Side(style="thin", color="B7B7B7"),
                )
        for col in sheet.columns:
            col_letter = col[0].column_letter
            sheet.column_dimensions[col_letter].width = 18 if col_letter in ["A", "B"] else 32
    out = MAIN / "02_Literature_Mapping" / "Literature_Mapping_Comparison_Matrix.xlsx"
    wb.save(out)
    copy_if_exists(out, REPORTS / "Literature_Mapping_Comparison_Matrix.xlsx")


def create_ieee_docx():
    doc = Document()
    doc.core_properties.author = STUDENT
    doc.core_properties.last_modified_by = STUDENT
    doc.core_properties.title = "NetGuard AI IEEE Draft"
    doc.core_properties.subject = "AI-based Network Anomaly Detection"
    doc.core_properties.comments = ""
    section = doc.sections[0]
    section.top_margin = Inches(0.75)
    section.bottom_margin = Inches(0.75)
    section.left_margin = Inches(0.7)
    section.right_margin = Inches(0.7)
    set_doc_styles(doc)

    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run("NetGuard AI: A Lightweight Machine Learning-Based Network Anomaly Detection Dashboard for Educational Infrastructure")
    run.bold = True
    run.font.name = "Times New Roman"
    run.font.size = Pt(16)
    run.font.color.rgb = RGBColor(0, 0, 0)

    meta = doc.add_paragraph()
    meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    meta.add_run(f"{STUDENT} ({NIM})\n{PROGRAM}, {FACULTY}\n{UNIVERSITY}, {YEAR}").font.size = Pt(10)

    doc.add_heading("Abstract", level=1)
    doc.add_paragraph(
        "Network reliability and security are critical for educational institutions, laboratories, and small organizations. "
        "However, many beginner network administrators still rely on manual inspection or basic monitoring tools that do not translate anomaly detection results into operational guidance. "
        "This paper proposes NetGuard AI, a lightweight machine learning-based network anomaly detection system using CICIDS2017-style traffic data. "
        "The system performs preprocessing, binary label conversion, model training, model evaluation, prediction, risk classification, and web-based visualization. "
        "Three classical machine learning algorithms are compared: Logistic Regression, Decision Tree, and Random Forest. "
        "Performance is evaluated using Accuracy, Precision, Recall, F1-score, and Confusion Matrix. "
        "The best model is selected based on F1-score, followed by Recall and Accuracy. "
        "In addition to model evaluation, NetGuard AI provides a Flask dashboard that displays total records, normal traffic, anomaly traffic, model accuracy, risk score, prediction results, and recommended actions. "
        "The novelty of this work lies in integrating machine learning anomaly detection with a lightweight educational network monitoring dashboard and actionable recommendations for beginner network administrators. "
        "The current implementation demonstrates a reproducible and low-budget prototype; final research submission should use a larger real CICIDS2017 subset for stronger empirical validity."
    )
    doc.add_paragraph("Keywords: Network Anomaly Detection, Intrusion Detection System, Machine Learning, CICIDS2017, Flask Dashboard")

    sections = [
        ("I. Introduction", [
            "Artificial Intelligence has increasingly been applied to cyber security, especially in intrusion detection and network anomaly monitoring. Educational institutions require reliable network services, but often lack expensive monitoring infrastructure. This motivates a lightweight and reproducible system that can be implemented on a normal laptop.",
            "The research problem addressed in this work is how to build an AI-based anomaly detection prototype that is not limited to model metrics but also provides a dashboard, risk classification, and actionable recommendations. The proposed project is NetGuard AI, a student-oriented system for detecting normal and anomalous traffic from CICIDS2017-style data.",
            "The contributions are: (1) a CICIDS2017 preprocessing pipeline, (2) comparison of Logistic Regression, Decision Tree, and Random Forest, (3) a Flask dashboard for monitoring visualization, and (4) risk-level and recommended-action generation for beginner network administrators.",
        ]),
        ("II. Related Works", [
            "Recent IDS literature includes classical machine learning benchmarking, deep learning, hybrid ML/DL models, feature importance, and dataset-validity studies. Maseer et al. benchmarked machine learning algorithms on CICIDS2017, while Dube warned that faulty use of CICIDS2017 can lead to misleading conclusions. Other studies use hybrid CNN-LSTM, DBN, broad learning, oversampling, feature embedding, and fuzzy clustering.",
            "Although many studies report high accuracy, several limitations remain. First, high performance on a single dataset may not generalize to operational networks. Second, complex deep learning models are not always suitable for low-budget educational implementation. Third, many papers stop at model evaluation and do not deliver a dashboard artifact for beginner administrators.",
        ]),
        ("III. Research Gap and Novelty", [
            "The primary research gap is the limited availability of reproducible, low-budget, dashboard-integrated machine learning anomaly detection systems for educational network infrastructure. Existing works often prioritize model performance, whereas NetGuard AI also emphasizes deployment-ready artifacts, dashboard visualization, risk classification, and user-oriented recommendations.",
            "The novelty is the integration of machine learning anomaly detection with a lightweight educational network monitoring dashboard and automated risk classification. This novelty is realistic, measurable, and implementable within an undergraduate UAS project.",
        ]),
        ("IV. Proposed Method", [
            "The proposed workflow consists of dataset collection, preprocessing, feature handling, model training, evaluation, best-model selection, prediction, risk classification, and dashboard visualization. CICIDS2017 labels are converted into binary classes: 0 for Normal and 1 for Anomaly. Missing values and infinite values are handled during preprocessing.",
            "The baseline models are Logistic Regression and Decision Tree. The proposed primary model is Random Forest because ensemble tree-based models are often strong for tabular network-flow data. Nevertheless, the final best model is selected objectively using F1-score, Recall, and Accuracy.",
        ]),
        ("V. Experimental Setup", [
            "The experiment uses Python, Pandas, NumPy, Scikit-Learn, Joblib, Matplotlib, Flask, Bootstrap, and Chart.js. The data split uses 80% training and 20% testing with random_state = 42. Logistic Regression uses StandardScaler and max_iter = 1000. Decision Tree uses random_state = 42. Random Forest uses n_estimators = 100 and random_state = 42.",
            "Evaluation metrics include Accuracy, Precision, Recall, F1-score, and Confusion Matrix. Final UAS submission must replace sample results with real CICIDS2017 subset results.",
        ]),
    ]
    for heading, paras in sections:
        doc.add_heading(heading, level=1)
        for para in paras:
            p = doc.add_paragraph(para)
            p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY

    doc.add_heading("VI. Results and Discussion", level=1)
    doc.add_paragraph(
        "Table I shows the current smoke-test result generated from the small sample CSV. The asterisk indicates that the result is only for pipeline validation and must not be used as the final research claim."
    )
    add_doc_table(doc, ["Model", "Accuracy", "Precision", "Recall", "F1-score", "Role"], EXPERIMENT_ROWS, [1.3, 0.9, 0.9, 0.9, 0.9, 1.4])
    doc.add_paragraph(
        "The current sample result reaches 1.0000 across all metrics because the sample dataset is intentionally small and separable. This is useful to prove that preprocessing, training, evaluation, prediction, export, and dashboard workflows run correctly. However, a final scientific conclusion requires a larger CICIDS2017 subset. In the final analysis, the student must discuss whether Random Forest outperforms the baselines, how class imbalance affects Recall and F1-score, and whether false positives or false negatives create operational risk."
    )

    doc.add_heading("VII. Conclusion", level=1)
    doc.add_paragraph(
        "NetGuard AI demonstrates a low-budget machine learning anomaly detection prototype with dashboard visualization and research artifacts. The project satisfies the implementation requirement of the UAS by providing preprocessing, baseline comparison, model training, evaluation, prediction, and web dashboard. Future work should include real CICIDS2017 experiments, cross-validation, hyperparameter tuning, feature importance, multi-dataset validation, and real-time network-log ingestion."
    )

    doc.add_heading("References", level=1)
    refs = build_references()
    for ref in refs:
        doc.add_paragraph(ref)

    out = MAIN / "09_Draft_IEEE" / "Draft_Artikel_IEEE.docx"
    doc.save(out)
    copy_if_exists(out, ROOT / "Draft_Artikel_IEEE_NetGuard_AI_Achmad_Maulana.docx")
    copy_if_exists(out, REPORTS / "Draft_Artikel_IEEE_NetGuard_AI_Achmad_Maulana.docx")


def build_references():
    refs = []
    for i, p in enumerate(PAPERS, 1):
        refs.append(
            f"[{i}] {p['authors']}, \"{p['title']},\" {p['publisher']}, {p['year']}, doi: {p['doi']}."
        )
    refs.extend([
        "[11] I. Sharafaldin, A. H. Lashkari, and A. A. Ghorbani, \"Toward generating a new intrusion detection dataset and intrusion traffic characterization,\" ICISSP, 2018.",
        "[12] S. Wang et al., \"Machine learning in network anomaly detection: A survey,\" IEEE Access, vol. 9, pp. 152379-152396, 2021.",
        "[13] G. Kocher and G. Kumar, \"Machine learning and deep learning methods for intrusion detection systems: recent developments and challenges,\" Soft Computing, vol. 25, pp. 9731-9763, 2021.",
        "[14] M. Tavallaee, E. Bagheri, W. Lu, and A. A. Ghorbani, \"A detailed analysis of the KDD CUP 99 data set,\" IEEE CISDA, 2009.",
        "[15] N. Moustafa and J. Slay, \"UNSW-NB15: a comprehensive data set for network intrusion detection systems,\" MilCIS, 2015.",
        "[16] Canadian Institute for Cybersecurity, \"CICIDS2017 Dataset,\" University of New Brunswick. [Online]. Available: https://www.unb.ca/cic/datasets/ids-2017.html.",
        "[17] F. Pedregosa et al., \"Scikit-learn: Machine learning in Python,\" Journal of Machine Learning Research, vol. 12, pp. 2825-2830, 2011.",
        "[18] W. McKinney, \"Data structures for statistical computing in Python,\" Proceedings of SciPy, 2010.",
        "[19] L. Breiman, \"Random forests,\" Machine Learning, vol. 45, pp. 5-32, 2001.",
        "[20] J. R. Quinlan, \"Induction of decision trees,\" Machine Learning, vol. 1, pp. 81-106, 1986.",
    ])
    return refs


def create_pdf_from_story(path, title, sections, landscape_mode=False):
    page_size = landscape(letter) if landscape_mode else letter
    doc = SimpleDocTemplate(str(path), pagesize=page_size, rightMargin=0.55*inch, leftMargin=0.55*inch, topMargin=0.55*inch, bottomMargin=0.55*inch)
    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle("TitleCenter", parent=styles["Title"], fontName="Times-Bold", fontSize=16, alignment=TA_CENTER, textColor=colors.black, spaceAfter=12))
    styles.add(ParagraphStyle("H", parent=styles["Heading2"], fontName="Times-Bold", fontSize=12, textColor=colors.black, spaceBefore=8, spaceAfter=4))
    styles.add(ParagraphStyle("BodyJ", parent=styles["BodyText"], fontName="Times-Roman", fontSize=9, leading=12, alignment=TA_JUSTIFY, textColor=colors.black))
    story = [Paragraph(title, styles["TitleCenter"]), Paragraph(f"{STUDENT} - {NIM}<br/>{PROGRAM}, {UNIVERSITY}", styles["BodyJ"]), Spacer(1, 8)]
    for heading, body in sections:
        story.append(Paragraph(heading, styles["H"]))
        if isinstance(body, list):
            for b in body:
                story.append(Paragraph(b, styles["BodyJ"]))
                story.append(Spacer(1, 4))
        else:
            story.append(body)
            story.append(Spacer(1, 8))
    doc.build(story)


def create_ieee_pdf():
    sections = [
        ("Abstract", [
            "Network reliability and security are critical for educational infrastructure. NetGuard AI proposes a lightweight machine learning anomaly detection system using CICIDS2017-style data, classical ML baselines, Random Forest as the proposed model, and a Flask dashboard with risk classification and recommended actions."
        ]),
        ("Introduction", [
            "The project addresses a practical gap: many IDS studies stop at model metrics, while beginner administrators need dashboard interpretation and operational guidance."
        ]),
        ("Method", [
            "The workflow includes preprocessing, binary label conversion, train/test split, Logistic Regression, Decision Tree, Random Forest, model selection, prediction, and dashboard visualization."
        ]),
        ("Results", [
            "Current sample results are smoke-test artifacts only. Final submission must use a real CICIDS2017 subset and replace all values with real metrics."
        ]),
        ("Conclusion", [
            "NetGuard AI is a reproducible educational prototype and can be improved with cross-validation, hyperparameter tuning, feature importance, and live network ingestion."
        ]),
        ("References", build_references()),
    ]
    out = MAIN / "09_Draft_IEEE" / "Draft_Artikel_IEEE.pdf"
    create_pdf_from_story(out, "Draft IEEE Article - NetGuard AI", sections)
    copy_if_exists(out, ROOT / "Draft_Artikel_IEEE_NetGuard_AI_Achmad_Maulana.pdf")
    copy_if_exists(out, REPORTS / "Draft_Artikel_IEEE_NetGuard_AI_Achmad_Maulana.pdf")


def add_slide_title(slide, title, subtitle=None):
    shape = slide.shapes.add_textbox(PptInches(0.55), PptInches(0.32), PptInches(12.2), PptInches(0.7))
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Aptos Display"
    p.font.size = PptPt(28)
    p.font.bold = True
    p.font.color.rgb = PptRGB(255, 255, 255)
    if subtitle:
        sub = slide.shapes.add_textbox(PptInches(0.58), PptInches(1.02), PptInches(11.7), PptInches(0.35))
        p2 = sub.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.size = PptPt(11)
        p2.font.color.rgb = PptRGB(195, 212, 235)


def add_pill(slide, x, y, w, h, text, fill, font=12):
    box = slide.shapes.add_shape(5, PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = PptRGB(70, 95, 130)
    tf = box.text_frame
    tf.text = text
    tf.paragraphs[0].font.size = PptPt(font)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = PptRGB(255, 255, 255)
    tf.paragraphs[0].alignment = 1
    return box


def add_body(slide, bullets, x=0.75, y=1.45, w=5.8, h=4.8, font=15):
    tb = slide.shapes.add_textbox(PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    tf = tb.text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.font.size = PptPt(font)
        p.font.color.rgb = PptRGB(235, 242, 255)
        p.space_after = PptPt(7)
    return tb


def create_pptx():
    prs = Presentation()
    prs.core_properties.author = STUDENT
    prs.core_properties.title = "NetGuard AI UAS Defense"
    prs.core_properties.subject = "AI-based Network Anomaly Detection"
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)
    bg_colors = [
        (9, 16, 34), (16, 35, 62), (24, 30, 62), (11, 44, 51),
        (30, 26, 61), (38, 38, 38), (13, 45, 80), (18, 58, 49),
        (53, 37, 72), (16, 16, 28), (10, 35, 58), (20, 20, 20),
    ]
    slides = [
        ("NetGuard AI", ["Predictive Network Failure & Anomaly Monitoring", STUDENT, NIM, UNIVERSITY], "cover"),
        ("Problem", ["Gangguan jaringan sering diketahui setelah layanan turun.", "Monitoring tradisional belum menerjemahkan data menjadi risiko.", "Admin pemula butuh dashboard yang mudah dipahami."], "problem"),
        ("Research Gap", ["Banyak studi IDS berhenti di akurasi model.", "Risiko faulty use CICIDS2017 masih perlu diantisipasi.", "Deployment edukatif dan rekomendasi tindakan masih jarang dibahas."], "gap"),
        ("Novelty", ["Main novelty: ML anomaly detection + lightweight educational dashboard.", "Secondary novelty: automated risk level + actionable recommendations.", "Konteks: sekolah, kampus, lab, dan UMKM low-budget."], "novelty"),
        ("Research Method", ["RM1: Dashboard-based ML anomaly detection.", "RM2: Baseline Logistic Regression/Decision Tree vs Random Forest.", "Dataset utama: CICIDS2017-style traffic flow CSV."], "rm"),
        ("AI Pipeline", ["Dataset", "Preprocessing", "Feature Handling", "Training", "Evaluation", "Prediction", "Dashboard"], "pipeline"),
        ("Experiment Design", ["Split 80:20, random_state 42.", "Metrics: Accuracy, Precision, Recall, F1-score, Confusion Matrix.", "Best model rule: F1-score, Recall, Accuracy."], "experiment"),
        ("Current Result", ["Sample CSV only validates pipeline.", "All models show 1.0000 because sample is tiny and simple.", "Final research must use real CICIDS2017 subset."], "result"),
        ("Dashboard Demo", ["Cards: total, normal, anomaly, accuracy.", "Charts: model comparison and confusion matrix.", "Pages: upload, train, predict, report."], "dashboard"),
        ("Deployment Plan", ["Frontend future: Next.js + Tailwind + Framer Motion.", "Backend current: Flask API; scalable future: NestJS/PostgreSQL.", "Deploy: Vercel + Railway/Render + GitHub Actions."], "deploy"),
        ("Defense Message", ["This is not only a literature review.", "It includes working preprocessing, training, evaluation, prediction, dashboard, and artifacts.", "Scientific strength depends on final real CICIDS2017 experiment."], "defense"),
        ("Conclusion", ["NetGuard AI is a low-budget AI cyber-security prototype.", "It is suitable for UAS, final project, and thesis development.", "Next: real dataset, tuning, cross-validation, feature importance."], "conclusion"),
    ]
    for idx, (title, bullets, kind) in enumerate(slides):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        r, g, b = bg_colors[idx % len(bg_colors)]
        rect = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        rect.fill.solid()
        rect.fill.fore_color.rgb = PptRGB(r, g, b)
        rect.line.fill.background()
        add_slide_title(slide, title, "NetGuard AI UAS Artificial Intelligence Defense Package" if idx else None)
        if kind == "cover":
            add_body(slide, bullets, x=0.8, y=1.7, w=8.0, h=3.0, font=22)
            add_pill(slide, 8.8, 1.8, 3.5, 0.55, "AI + Cyber Security", PptRGB(30, 144, 255), 15)
            add_pill(slide, 8.8, 2.55, 3.5, 0.55, "CICIDS2017", PptRGB(46, 204, 113), 15)
            add_pill(slide, 8.8, 3.3, 3.5, 0.55, "Flask Dashboard", PptRGB(155, 89, 182), 15)
        elif kind == "pipeline":
            x = 0.7
            colors_list = [PptRGB(37, 99, 235), PptRGB(20, 184, 166), PptRGB(34, 197, 94), PptRGB(250, 204, 21), PptRGB(249, 115, 22), PptRGB(236, 72, 153), PptRGB(139, 92, 246)]
            for i, item in enumerate(bullets):
                add_pill(slide, x + (i % 4) * 3.05, 1.8 + (i // 4) * 1.35, 2.65, 0.75, item, colors_list[i], 14)
            add_body(slide, ["Narasi: pipeline ini menunjukkan bahwa proyek memiliki implementasi penuh, bukan hanya rangkuman jurnal."], x=0.9, y=5.0, w=11.0, h=1.0, font=16)
        else:
            add_body(slide, bullets, x=0.8, y=1.55, w=7.0, h=4.8, font=18)
            add_pill(slide, 8.6, 1.55, 3.6, 0.62, "Research Artifact", PptRGB(37, 99, 235), 14)
            add_pill(slide, 8.6, 2.35, 3.6, 0.62, "Reproducible", PptRGB(20, 184, 166), 14)
            add_pill(slide, 8.6, 3.15, 3.6, 0.62, "Low-Budget", PptRGB(34, 197, 94), 14)
            if kind == "dashboard":
                img = ROOT / "reports" / "figures" / "confusion_matrix.png"
                if img.exists():
                    slide.shapes.add_picture(str(img), PptInches(8.7), PptInches(4.0), width=PptInches(2.7))
        foot = slide.shapes.add_textbox(PptInches(0.6), PptInches(7.0), PptInches(12.0), PptInches(0.3))
        p = foot.text_frame.paragraphs[0]
        p.text = f"{STUDENT} | {NIM} | NetGuard AI | {idx+1}/12"
        p.font.size = PptPt(9)
        p.font.color.rgb = PptRGB(170, 185, 210)
    out = MAIN / "10_Presentasi" / "Slide_Presentasi.pptx"
    prs.save(out)
    copy_if_exists(out, ROOT / "Slide_Presentasi_NetGuard_AI_Achmad_Maulana.pptx")
    copy_if_exists(out, REPORTS / "Slide_Presentasi_NetGuard_AI_Achmad_Maulana.pptx")


def create_slide_pdf():
    slides = [
        ("NetGuard AI", "AI-based network anomaly detection with dashboard monitoring."),
        ("Problem", "Network failures are often detected late; beginner admins need simple risk-oriented tools."),
        ("Research Gap", "Prior IDS works often focus on metrics, not educational dashboard deployment."),
        ("Novelty", "ML anomaly detection integrated with lightweight dashboard and recommendations."),
        ("Research Method", "Compare Logistic Regression, Decision Tree, and Random Forest on CICIDS2017."),
        ("Pipeline", "Dataset, preprocessing, training, evaluation, prediction, dashboard."),
        ("Experiment", "80/20 split, seed 42, Accuracy, Precision, Recall, F1, Confusion Matrix."),
        ("Result", "Sample result validates pipeline; final claim requires real CICIDS2017 subset."),
        ("Dashboard", "Cards, charts, prediction, report, and API summary."),
        ("Deployment", "Current Flask; future Next.js, Tailwind, Framer Motion, PostgreSQL."),
        ("Defense", "Implementation exists; not only literature review."),
        ("Conclusion", "Ready as UAS prototype; strengthen with real data and tuning."),
    ]
    path = MAIN / "10_Presentasi" / "Slide_Presentasi.pdf"
    doc = SimpleDocTemplate(str(path), pagesize=landscape(letter), rightMargin=0.5*inch, leftMargin=0.5*inch, topMargin=0.5*inch, bottomMargin=0.5*inch)
    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle("SlideTitle", parent=styles["Title"], fontName="Helvetica-Bold", fontSize=28, alignment=TA_CENTER, textColor=colors.HexColor("#0B1020")))
    styles.add(ParagraphStyle("SlideBody", parent=styles["BodyText"], fontName="Helvetica", fontSize=18, leading=26, alignment=TA_CENTER, textColor=colors.HexColor("#1F2937")))
    story = []
    for i, (t, b) in enumerate(slides):
        story.append(Spacer(1, 1.2*inch))
        story.append(Paragraph(t, styles["SlideTitle"]))
        story.append(Spacer(1, 0.35*inch))
        story.append(Paragraph(b, styles["SlideBody"]))
        story.append(Spacer(1, 0.5*inch))
        story.append(Paragraph(f"{STUDENT} | {NIM} | Slide {i+1}/12", styles["SlideBody"]))
        if i != len(slides) - 1:
            story.append(PageBreak())
    doc.build(story)
    copy_if_exists(path, ROOT / "Slide_Presentasi_NetGuard_AI_Achmad_Maulana.pdf")


def create_mapping_markdown():
    out = MAIN / "02_Literature_Mapping" / "Literature_Mapping.md"
    lines = ["# Literature Mapping - NetGuard AI", "", f"Mahasiswa: {STUDENT} ({NIM})", ""]
    lines.append("| No | Tahun | Penulis | Dataset | Metode | Hasil | Kelebihan | Kekurangan |")
    lines.append("|---:|---:|---|---|---|---|---|---|")
    for p in PAPERS:
        lines.append(f"| {p['no']} | {p['year']} | {p['authors']} | {p['dataset']} | {p['method']} | {p['result']} | {p['strength']} | {p['weakness']} |")
    lines.append("\n## Primary Research Gap\n")
    lines.append("Belum banyak penelitian deteksi anomali jaringan berbasis CICIDS2017 yang mengintegrasikan model machine learning sederhana, dashboard monitoring ringan, klasifikasi risiko otomatis, rekomendasi tindakan, dan dokumentasi reproducible untuk konteks pendidikan low-budget.")
    out.write_text("\n".join(lines), encoding="utf-8")
    copy_if_exists(out, REPORTS / "Literature_Mapping_Final.md")


def create_gap_docs():
    text = f"""# Gap Analysis, Novelty, and Research Method

Mahasiswa: {STUDENT} ({NIM})

## Research Gap

Primary research gap: belum banyak penelitian deteksi anomali jaringan berbasis CICIDS2017 yang mengintegrasikan model machine learning sederhana, dashboard monitoring ringan, klasifikasi risiko otomatis, rekomendasi tindakan, dan dokumentasi reproducible untuk konteks pendidikan low-budget.

## Evidence

1. Studi benchmarking seperti Maseer et al. membahas performa model, tetapi bukan dashboard edukatif.
2. Dube menekankan risiko faulty use CICIDS2017, sehingga pipeline harus transparan.
3. Studi hybrid ML/DL dan deep learning menunjukkan performa tinggi, tetapi kompleksitasnya lebih tinggi untuk laptop mahasiswa.
4. Studi Random Forest dan feature importance relevan, tetapi belum berfokus pada pengalaman admin pemula.

## Novelty Statement

NetGuard AI mengintegrasikan machine learning anomaly detection dengan lightweight educational network monitoring dashboard, automated risk classification, dan actionable recommendation untuk administrator jaringan pemula.

## Research Method

RM1: Implementasi NetGuard AI Dashboard-Based ML Anomaly Detection.

RM2: Perbandingan baseline Logistic Regression dan Decision Tree dengan proposed Random Forest pada dataset CICIDS2017.

## Proposed Method

Metode usulan adalah Random Forest sebagai model utama dalam kerangka dashboard-based anomaly monitoring. Logistic Regression dan Decision Tree digunakan sebagai baseline pembanding.
"""
    out = MAIN / "03_Gap_Analysis" / "Gap_Novelty_RM.md"
    out.write_text(text, encoding="utf-8")
    copy_if_exists(out, REPORTS / "Gap_Novelty_RM_Final.md")


def create_model_config_and_logs():
    model_config = {
        "student": STUDENT,
        "nim": NIM,
        "project": "NetGuard AI",
        "dataset": "CICIDS2017-style sample; final submission must use real CICIDS2017 subset",
        "models": {
            "Logistic Regression": {"max_iter": 1000, "random_state": 42, "scaler": "StandardScaler"},
            "Decision Tree": {"random_state": 42},
            "Random Forest": {"n_estimators": 100, "random_state": 42},
        },
        "selection_rule": "F1-score, then Recall, then Accuracy",
        "created": str(date.today()),
    }
    out = MAIN / "06_Model" / "model_configuration.json"
    out.write_text(json.dumps(model_config, indent=4), encoding="utf-8")
    log = MAIN / "07_Hasil_Eksperimen" / "Experiment_Log.md"
    log.write_text(
        "# Experiment Log\n\n"
        "Command flow:\n\n"
        "```cmd\n"
        "py -3 -m pip install -r requirements.txt\n"
        "py -3 app\\preprocessing.py\n"
        "py -3 app\\train.py\n"
        "py -3 app\\test.py\n"
        "py -3 app\\predict.py\n"
        "py -3 app\\main.py\n"
        "```\n\n"
        "Current sample result validates pipeline only. Final UAS result must use real CICIDS2017 subset.\n",
        encoding="utf-8",
    )


def create_turnitin_readiness_pdf():
    sections = [
        ("Turnitin Requirement", [
            "Similarity must be maximum 15%. Single-source similarity must be maximum 3%. References are excluded from similarity checking if the institution setting supports exclusion."
        ]),
        ("Revision Strategy", [
            "Paraphrase every literature paragraph, avoid copying abstracts, cite IEEE style, use original analysis for gap and novelty, and check similarity before final upload."
        ]),
        ("Manual Step", [
            "The real Turnitin_Report.pdf must be generated from the paid Turnitin system by the student or institution. This file is a readiness guide, not a substitute for the official Turnitin report."
        ]),
    ]
    create_pdf_from_story(MAIN / "11_Turnitin" / "Turnitin_Readiness_Guide.pdf", "Turnitin Readiness Guide", sections)


def create_deployment_source():
    dep = MAIN / "12_Deployment"
    project = dep / "netguard-next-platform"
    (project / "src/app").mkdir(parents=True, exist_ok=True)
    (project / "src/components").mkdir(parents=True, exist_ok=True)
    (project / "src/lib").mkdir(parents=True, exist_ok=True)
    (project / "prisma").mkdir(parents=True, exist_ok=True)
    (project / ".github/workflows").mkdir(parents=True, exist_ok=True)
    files = {
        project / "package.json": json.dumps({
            "name": "netguard-ai-platform",
            "version": "1.0.0",
            "private": True,
            "scripts": {"dev": "next dev", "build": "next build", "start": "next start", "lint": "next lint"},
            "dependencies": {
                "next": "latest",
                "react": "latest",
                "react-dom": "latest",
                "framer-motion": "latest",
                "three": "latest",
                "@react-three/fiber": "latest",
                "tailwindcss": "latest",
                "lucide-react": "latest",
                "@prisma/client": "latest",
                "zod": "latest",
            },
            "devDependencies": {"typescript": "latest", "prisma": "latest", "@types/node": "latest", "@types/react": "latest"}
        }, indent=2),
        project / "README.md": """# NetGuard AI Next.js Deployment Concept

This is a deployment-ready concept source for a future world-class UI.

Recommended stack:
- Frontend: Next.js, React, Tailwind CSS, Framer Motion, Three.js
- Backend API: Next.js route handlers or NestJS
- Database: PostgreSQL with Prisma
- Deployment: Vercel for frontend, Railway/Render for PostgreSQL/API

Run:
```cmd
npm install
npm run dev
```
""",
        project / "src/app/page.tsx": """import { MotionDashboard } from '../components/MotionDashboard';

export default function Home() {
  return <MotionDashboard />;
}
""",
        project / "src/components/MotionDashboard.tsx": """'use client';
import { motion, useScroll, useTransform } from 'framer-motion';

export function MotionDashboard() {
  const { scrollYProgress } = useScroll();
  const background = useTransform(scrollYProgress, [0, 0.5, 1], ['#0b1020', '#12323f', '#201b3f']);
  return (
    <motion.main style={{ background }} className="min-h-screen text-white overflow-hidden">
      <CursorGlow />
      <section className="h-screen flex items-center px-10">
        <motion.div initial={{ opacity: 0, y: 40 }} animate={{ opacity: 1, y: 0 }} transition={{ duration: .8 }}>
          <p className="uppercase tracking-widest text-cyan-300">NetGuard AI</p>
          <h1 className="text-7xl font-semibold max-w-5xl">Predictive Network Failure & Anomaly Monitoring</h1>
          <p className="mt-6 text-xl text-slate-300 max-w-3xl">Apple-clean, Hugging Face-like processing states, and Microsoft-grade operational clarity.</p>
        </motion.div>
      </section>
      <section className="min-h-screen sticky top-0 grid place-items-center">
        <motion.div className="rounded-3xl border border-cyan-300/30 bg-white/10 p-8 backdrop-blur-xl shadow-2xl">
          <div className="animate-pulse h-2 rounded-full bg-cyan-300 mb-6" />
          <h2 className="text-4xl font-bold">AI Processing Pipeline</h2>
          <p className="mt-4 text-slate-200">Upload -> Preprocess -> Train -> Predict -> Risk Action</p>
        </motion.div>
      </section>
    </motion.main>
  );
}

function CursorGlow() {
  return <div className="pointer-events-none fixed inset-0 bg-[radial-gradient(circle_at_var(--x,50%)_var(--y,50%),rgba(56,189,248,.22),transparent_25%)]" />;
}
""",
        project / "src/app/api/summary/route.ts": """import { NextResponse } from 'next/server';

export async function GET() {
  return NextResponse.json({
    totalRecords: 0,
    normalTraffic: 0,
    anomalyTraffic: 0,
    riskLevel: 'Low',
    message: 'Connect this endpoint to Flask model outputs or PostgreSQL records.'
  });
}
""",
        project / "prisma/schema.prisma": """generator client {
  provider = "prisma-client-js"
}

datasource db {
  provider = "postgresql"
  url      = env("DATABASE_URL")
}

model Experiment {
  id        String   @id @default(cuid())
  dataset   String
  modelName String
  accuracy  Float
  precision Float
  recall    Float
  f1Score   Float
  createdAt DateTime @default(now())

  predictions Prediction[]
}

model Prediction {
  id           String   @id @default(cuid())
  experimentId String
  label        Int
  labelText    String
  riskLevel    String
  createdAt    DateTime @default(now())

  experiment Experiment @relation(fields: [experimentId], references: [id])
  @@index([experimentId])
  @@index([riskLevel])
}
""",
        project / ".github/workflows/deploy.yml": """name: Deploy

on:
  push:
    branches: [ main ]

jobs:
  build:
    runs-on: ubuntu-latest
    steps:
      - uses: actions/checkout@v4
      - uses: actions/setup-node@v4
        with:
          node-version: 20
      - run: npm install
        working-directory: Achmad_Maulana_241730016_UAS_AI/12_Deployment/netguard-next-platform
      - run: npm run build
        working-directory: Achmad_Maulana_241730016_UAS_AI/12_Deployment/netguard-next-platform
""",
        dep / "Deployment_Documentation.md": """# Deployment Documentation

## Recommended Production Strategy

1. Current demo: Flask local dashboard.
2. Public deployment option: Render or Railway for Flask.
3. Premium future UI: Next.js + Tailwind + Framer Motion on Vercel.
4. Database: PostgreSQL on Railway/Supabase.
5. CI/CD: GitHub Actions.

## Evidence To Collect

- URL aplikasi setelah deploy.
- Screenshot dashboard.
- Screenshot prediction page.
- Screenshot report page.
- GitHub repository link.
- Deployment log.

## Limitation

Public URL cannot be created automatically without the student's Vercel/Render/Railway account.
""",
    }
    for path, content in files.items():
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_text(content, encoding="utf-8")


def copy_project_artifacts():
    copy_if_exists(ROOT / "data" / "sample" / "sample_cicids2017.csv", MAIN / "04_Dataset" / "Raw_Dataset" / "sample_cicids2017.csv")
    copy_if_exists(ROOT / "data" / "processed" / "sample_cicids2017_processed.csv", MAIN / "04_Dataset" / "Processed_Dataset" / "sample_cicids2017_processed.csv")
    for fname in ["preprocessing.py", "train.py", "test.py", "predict.py", "predictor.py", "evaluate.py", "main.py", "utils.py"]:
        copy_if_exists(ROOT / "app" / fname, MAIN / "05_Source_Code" / "Script" / fname)
    copy_if_exists(ROOT / "README.md", MAIN / "05_Source_Code" / "README.md")
    copy_if_exists(ROOT / "requirements.txt", MAIN / "05_Source_Code" / "requirements.txt")
    copy_if_exists(ROOT / "models" / "netguard_best_model.pkl", MAIN / "06_Model" / "netguard_best_model.pkl")
    for fname in ["metrics.json", "model_comparison.csv", "prediction_result.csv", "research_summary.json"]:
        copy_if_exists(ROOT / "reports" / fname, MAIN / "07_Hasil_Eksperimen" / fname)
    figs = ROOT / "reports" / "figures"
    if figs.exists():
        for f in figs.glob("*"):
            if f.is_file():
                copy_if_exists(f, MAIN / "08_Visualisasi" / f.name)


def update_main_readme():
    text = f"""# UAS AI Final Package - NetGuard AI

## Identitas

- Nama Mahasiswa: {STUDENT}
- NIM: {NIM}
- Program Studi: {PROGRAM}
- Fakultas: {FACULTY}
- Universitas: {UNIVERSITY}
- Tahun: {YEAR}

## Judul Penelitian

NetGuard AI: Predictive Network Failure and Anomaly Monitoring Berbasis Machine Learning pada Dataset CICIDS2017

## Research Gap

Belum banyak penelitian deteksi anomali jaringan berbasis CICIDS2017 yang mengintegrasikan model machine learning sederhana, dashboard monitoring ringan, klasifikasi risiko otomatis, rekomendasi tindakan, dan dokumentasi reproducible untuk konteks pendidikan low-budget.

## Novelty

Integrasi machine learning anomaly detection dengan lightweight educational network monitoring dashboard, automated risk classification, dan actionable recommendations untuk administrator jaringan pemula.

## Research Method

RM1: Implementasi NetGuard AI Dashboard-Based ML Anomaly Detection.

RM2: Perbandingan baseline Logistic Regression dan Decision Tree dengan proposed Random Forest pada dataset CICIDS2017.

## Command Reproduksi

```cmd
py -3 -m pip install -r requirements.txt
py -3 app\\preprocessing.py
py -3 app\\train.py
py -3 app\\test.py
py -3 app\\predict.py
py -3 app\\main.py
```

## Struktur Folder

Folder ini mengikuti ketentuan Google Drive UAS:

01_Paper, 02_Literature_Mapping, 03_Gap_Analysis, 04_Dataset, 05_Source_Code, 06_Model, 07_Hasil_Eksperimen, 08_Visualisasi, 09_Draft_IEEE, 10_Presentasi, 11_Turnitin, 12_Deployment.

## Catatan Kejujuran Akademik

Sample CSV hanya digunakan untuk smoke test. Untuk final submission, jalankan dataset CICIDS2017 nyata atau subset representatif, lalu perbarui metrik, confusion matrix, narasi hasil, dan artikel IEEE.
"""
    (MAIN / "README.md").write_text(text, encoding="utf-8")


def create_docx_pdf_support_docs():
    # Presentation narration
    narration = MAIN / "10_Presentasi" / "Narasi_Presentasi_Detail.md"
    narration.write_text(
        f"""# Narasi Presentasi Detail

## Opening

Assalamualaikum warahmatullahi wabarakatuh. Perkenalkan saya {STUDENT}, NIM {NIM}, dari {PROGRAM}, {UNIVERSITY}. Pada kesempatan ini saya mempresentasikan proyek UAS Artificial Intelligence berjudul NetGuard AI.

## Problem

Permasalahan utama adalah gangguan jaringan sering diketahui setelah layanan turun. Dalam konteks sekolah, kampus, laboratorium, dan UMKM, administrator jaringan membutuhkan sistem sederhana yang dapat membantu membaca potensi anomali trafik.

## Solution

NetGuard AI menggabungkan machine learning dengan dashboard monitoring. Sistem membaca dataset CICIDS2017-style, melakukan preprocessing, melatih beberapa model, mengevaluasi performa, melakukan prediksi, menghitung risk score, dan memberi recommended action.

## Method

Metode yang digunakan adalah Logistic Regression dan Decision Tree sebagai baseline, serta Random Forest sebagai proposed model. Evaluasi menggunakan Accuracy, Precision, Recall, F1-score, dan Confusion Matrix.

## Result

Hasil saat ini pada sample CSV menunjukkan pipeline berjalan penuh. Namun, hasil final akademik harus memakai subset CICIDS2017 nyata agar valid sebagai penelitian.

## Closing

Kontribusi utama NetGuard AI adalah integrasi model AI, dashboard monitoring ringan, risk classification, dan actionable recommendation untuk administrator jaringan pemula.
""",
        encoding="utf-8",
    )


def main():
    ensure_dirs()
    create_literature_xlsx()
    create_mapping_markdown()
    create_gap_docs()
    create_ieee_docx()
    create_ieee_pdf()
    create_pptx()
    create_slide_pdf()
    create_model_config_and_logs()
    create_turnitin_readiness_pdf()
    create_deployment_source()
    copy_project_artifacts()
    update_main_readme()
    create_docx_pdf_support_docs()
    print("Generated final UAS artifacts in", MAIN)


if __name__ == "__main__":
    main()
