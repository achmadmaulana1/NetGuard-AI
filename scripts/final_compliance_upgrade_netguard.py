from __future__ import annotations

import csv
import json
import math
import subprocess
import zipfile
from pathlib import Path

import numpy as np
import pandas as pd
from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor
from PIL import Image, ImageDraw, ImageFilter, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches as PptInches, Pt as PptPt
from pypdf import PdfReader


ROOT = Path(__file__).resolve().parents[1]
MAIN = ROOT / "Achmad_Maulana_241730016_UAS_AI"
REPORTS = ROOT / "reports"
FIGURES = REPORTS / "figures"

STUDENT = "Achmad Maulana"
NIM = "241730016"
PROGRAM = "Program Studi Informatika"
FACULTY = "Fakultas Sains dan Teknologi"
UNIVERSITY = "Universitas Islam Negeri Sultan Maulana Hasanuddin Banten"
YEAR = "2026"
TITLE = "NetGuard AI: Predictive Network Failure and Anomaly Monitoring"

FONT = "Times New Roman"
BLACK = RGBColor(0, 0, 0)


PAPERS = [
    {
        "key": "engelen2021troubleshooting",
        "title": "Troubleshooting an Intrusion Detection Dataset: The CICIDS2017 Case Study",
        "authors": "Engelen, G., Rimmer, V., and Joosen, W.",
        "year": "2021",
        "venue": "IEEE Security & Privacy Workshops / dataset quality study",
        "category": "IEEE/Scopus placeholder",
        "dataset": "CICIDS2017",
        "method": "Dataset troubleshooting and validation",
        "result": "Identifies quality issues and cautions in CICIDS2017 usage.",
        "weakness": "Focuses on dataset analysis, not an educational dashboard.",
        "file_hint": "15_Engelen_2021",
    },
    {
        "key": "sajid2024cloud",
        "title": "Machine Learning for Cloud and Network Intrusion Detection",
        "authors": "Sajid et al.",
        "year": "2024",
        "venue": "Journal of Cloud Computing",
        "category": "Scopus placeholder",
        "dataset": "Network intrusion datasets",
        "method": "Machine learning comparison",
        "result": "Reports strong performance for classical classifiers.",
        "weakness": "Operational dashboard integration is limited.",
        "file_hint": "02_Sajid_2024",
    },
    {
        "key": "dube2024cicids",
        "title": "On Faulty Use of the CICIDS2017 Dataset in Network Intrusion Detection Research",
        "authors": "Dube et al.",
        "year": "2024",
        "venue": "Dataset validity study",
        "category": "Scopus placeholder",
        "dataset": "CICIDS2017",
        "method": "Critical dataset review",
        "result": "Highlights methodological risks when using CICIDS2017.",
        "weakness": "Does not propose a beginner-ready monitoring prototype.",
        "file_hint": "03_Dube_2024",
    },
    {
        "key": "talukder2024bigdata",
        "title": "Network Intrusion Detection Using Machine Learning and Feature Engineering",
        "authors": "Talukder et al.",
        "year": "2024",
        "venue": "Journal of Big Data",
        "category": "Scopus placeholder",
        "dataset": "CICIDS2017 / related intrusion datasets",
        "method": "Feature engineering and classifier evaluation",
        "result": "Shows that feature preparation strongly affects performance.",
        "weakness": "Deployment and reproducible student artefacts are not central.",
        "file_hint": "04_Talukder_2024",
    },
    {
        "key": "abdelaziz2024rfimportance",
        "title": "Random Forest Permutation Feature Importance for Intrusion Detection",
        "authors": "Abdelaziz et al.",
        "year": "2024",
        "venue": "Feature importance study",
        "category": "Scopus placeholder",
        "dataset": "CICIDS2017 / cybersecurity dataset",
        "method": "Random Forest and permutation importance",
        "result": "Demonstrates interpretability potential of tree-based models.",
        "weakness": "Does not translate interpretation into dashboard actions.",
        "file_hint": "10_Abdelaziz_2024",
    },
    {
        "key": "vitorino2024adversarial",
        "title": "Adversarial Robustness of Network Intrusion Detection Systems",
        "authors": "Vitorino et al.",
        "year": "2024",
        "venue": "arXiv preprint",
        "category": "Preprint placeholder",
        "dataset": "Network intrusion datasets",
        "method": "Adversarial robustness evaluation",
        "result": "Shows robustness as an important future validation issue.",
        "weakness": "Not focused on low-budget educational deployment.",
        "file_hint": "13_Vitorino_2024",
    },
    {
        "key": "ahmed2025ids",
        "title": "Intrusion Detection Using Machine Learning for Modern Networks",
        "authors": "Ahmed et al.",
        "year": "2025",
        "venue": "Scientific Reports",
        "category": "Scopus/Nature placeholder",
        "dataset": "Public IDS datasets",
        "method": "Machine learning / deep learning comparison",
        "result": "Reports high detection performance on benchmark datasets.",
        "weakness": "Dashboard and novice administrator guidance remain limited.",
        "file_hint": "08_Ahmed_2025",
    },
    {
        "key": "pai2025eurasip",
        "title": "Information Security Intrusion Detection Using Machine Learning",
        "authors": "Pai et al.",
        "year": "2025",
        "venue": "EURASIP Journal on Information Security",
        "category": "Scopus placeholder",
        "dataset": "IDS benchmark datasets",
        "method": "Machine learning comparison",
        "result": "Supports model comparison as a valid IDS experiment design.",
        "weakness": "Does not emphasize reproducible local-laptop implementation.",
        "file_hint": "09_Pai_2025",
    },
    {
        "key": "nature2025iotids",
        "title": "Smart Deep Learning for IoT Intrusion Detection",
        "authors": "Authors listed in source PDF",
        "year": "2025",
        "venue": "Scientific Reports / Nature portfolio",
        "category": "Scopus/Nature placeholder",
        "dataset": "IoT / IDS datasets",
        "method": "Deep learning-based IDS",
        "result": "Demonstrates the benefit of advanced models for IDS.",
        "weakness": "Requires stronger compute and is less beginner-friendly.",
        "file_hint": "12_Nature_2025",
    },
    {
        "key": "ijca2025rf",
        "title": "High Accuracy CICIDS2017 Intrusion Detection Using Random Forest",
        "authors": "Authors listed in source PDF",
        "year": "2025",
        "venue": "IJCA",
        "category": "Journal placeholder",
        "dataset": "CICIDS2017",
        "method": "Random Forest",
        "result": "Shows tree-based model strength on CICIDS2017.",
        "weakness": "May be limited to accuracy reporting without dashboard context.",
        "file_hint": "16_IJCA_2025",
    },
    {
        "key": "isa2025ensemble",
        "title": "Ensemble Learning for CICIDS2017 and CICIDS2018 Intrusion Detection",
        "authors": "Authors listed in source PDF",
        "year": "2025",
        "venue": "Information Security Applications / IDS study",
        "category": "Scopus placeholder",
        "dataset": "CICIDS2017 and CICIDS2018",
        "method": "Ensemble learning",
        "result": "Supports ensemble and cross-dataset validation direction.",
        "weakness": "More complex than the current low-budget teaching target.",
        "file_hint": "17_ISA_2025",
    },
]


def ensure_dirs() -> None:
    for rel in [
        "01_Paper",
        "02_Literature_Mapping",
        "03_Gap_Analysis",
        "04_Dataset/Raw_Dataset",
        "04_Dataset/Processed_Dataset",
        "05_Source_Code/Notebook",
        "05_Source_Code/Script",
        "06_Model",
        "07_Hasil_Eksperimen",
        "08_Visualisasi",
        "09_Draft_IEEE",
        "10_Presentasi",
        "11_Turnitin",
        "12_Deployment",
        "13_GitHub",
        "14_Dokumentasi",
        "15_Bukti_Submit",
    ]:
        (MAIN / rel).mkdir(parents=True, exist_ok=True)


def read_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8")) if path.exists() else {}


def set_doc_props(doc: Document, title: str = "") -> None:
    props = doc.core_properties
    props.author = STUDENT
    props.last_modified_by = STUDENT
    props.comments = ""
    props.subject = ""
    props.keywords = ""
    if title:
        props.title = title


def style_doc(doc: Document, body_size: int = 10) -> None:
    for style_name in ["Normal", "Heading 1", "Heading 2", "Heading 3", "Title", "Subtitle"]:
        if style_name in doc.styles:
            style = doc.styles[style_name]
            style.font.name = FONT
            style._element.rPr.rFonts.set(qn("w:eastAsia"), FONT)
            style.font.color.rgb = BLACK
            if style_name == "Normal":
                style.font.size = Pt(body_size)
    for p in doc.paragraphs:
        for r in p.runs:
            r.font.name = FONT
            r._element.rPr.rFonts.set(qn("w:eastAsia"), FONT)
            r.font.color.rgb = BLACK


def add_paragraph(doc: Document, text: str, size: int = 10, bold: bool = False, align=None, italic: bool = False) -> None:
    p = doc.add_paragraph()
    if align is not None:
        p.alignment = align
    p.paragraph_format.space_after = Pt(4)
    r = p.add_run(text)
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.color.rgb = BLACK
    r.bold = bold
    r.italic = italic


def add_heading(doc: Document, text: str, level: int = 1) -> None:
    p = doc.add_paragraph()
    p.paragraph_format.space_before = Pt(7)
    p.paragraph_format.space_after = Pt(4)
    r = p.add_run(text)
    r.bold = True
    r.font.name = FONT
    r.font.size = Pt(12 if level == 1 else 10)
    r.font.color.rgb = BLACK


def table_cell(cell, text: str, bold: bool = False, size: int = 8, align=WD_ALIGN_PARAGRAPH.CENTER) -> None:
    cell.text = ""
    p = cell.paragraphs[0]
    p.alignment = align
    p.paragraph_format.space_after = Pt(0)
    r = p.add_run(str(text))
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.color.rgb = BLACK
    r.bold = bold
    cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER


def add_table(doc: Document, headers: list[str], rows: list[list[str]], caption: str, font_size: int = 7) -> None:
    add_paragraph(doc, caption, size=8, bold=True, align=WD_ALIGN_PARAGRAPH.CENTER)
    table = doc.add_table(rows=1, cols=len(headers))
    table.style = "Table Grid"
    for i, h in enumerate(headers):
        table_cell(table.rows[0].cells[i], h, True, font_size)
    for row in rows:
        cells = table.add_row().cells
        for i, value in enumerate(row):
            table_cell(cells[i], value, False, font_size, WD_ALIGN_PARAGRAPH.LEFT if len(str(value)) > 20 else WD_ALIGN_PARAGRAPH.CENTER)


def add_two_columns(section) -> None:
    cols = section._sectPr.xpath("./w:cols")
    cols = cols[0] if cols else OxmlElement("w:cols")
    if cols.getparent() is None:
        section._sectPr.append(cols)
    cols.set(qn("w:num"), "2")
    cols.set(qn("w:space"), "360")


def load_metrics_and_data() -> tuple[dict, pd.DataFrame, pd.DataFrame]:
    metrics = read_json(REPORTS / "metrics.json")
    comparison = pd.read_csv(REPORTS / "model_comparison.csv")
    data_path = ROOT / "data" / "processed" / "cicids2017_ddos_representative_20000_processed.csv"
    df = pd.read_csv(data_path) if data_path.exists() else pd.DataFrame()
    return metrics, comparison, df


def font(size: int = 24) -> ImageFont.ImageFont:
    for candidate in [
        "C:/Windows/Fonts/times.ttf",
        "C:/Windows/Fonts/arial.ttf",
        "C:/Windows/Fonts/calibri.ttf",
    ]:
        if Path(candidate).exists():
            return ImageFont.truetype(candidate, size)
    return ImageFont.load_default()


def draw_background(path: Path, w: int = 1920, h: int = 1080) -> None:
    img = Image.new("RGB", (w, h), "#eef6ff")
    px = img.load()
    for y in range(h):
        for x in range(w):
            t = x / w
            u = y / h
            r = int(232 - 40 * u + 10 * math.sin(t * 5))
            g = int(242 - 25 * t)
            b = int(255 - 30 * u + 30 * t)
            px[x, y] = (max(0, r), max(0, g), max(0, b))
    layer = Image.new("RGBA", (w, h), (0, 0, 0, 0))
    d = ImageDraw.Draw(layer)
    blobs = [
        (100, 40, 850, 620, "#7dd3fc", 95),
        (920, 120, 1800, 900, "#a7f3d0", 90),
        (520, 580, 1450, 1180, "#c4b5fd", 75),
    ]
    for x1, y1, x2, y2, color, alpha in blobs:
        d.ellipse((x1, y1, x2, y2), fill=(*Image.new("RGB", (1, 1), color).getpixel((0, 0)), alpha))
    layer = layer.filter(ImageFilter.GaussianBlur(90))
    img = Image.alpha_composite(img.convert("RGBA"), layer)
    d = ImageDraw.Draw(img)
    rng = np.random.default_rng(42)
    pts = [(int(rng.integers(90, w - 90)), int(rng.integers(90, h - 90))) for _ in range(55)]
    for i, p in enumerate(pts):
        for q in pts[i + 1 :]:
            dist = math.dist(p, q)
            if dist < 245:
                alpha = int(78 * (1 - dist / 245))
                d.line([p, q], fill=(27, 75, 120, alpha), width=2)
    for p in pts:
        d.ellipse((p[0] - 5, p[1] - 5, p[0] + 5, p[1] + 5), fill=(14, 116, 144, 150))
    d.rounded_rectangle((118, 104, 1800, 970), radius=42, outline=(255, 255, 255, 170), width=3)
    img.convert("RGB").save(path, quality=95)


def draw_bar_chart(path: Path, title: str, labels: list[str], series: dict[str, list[float]], percent: bool = True) -> None:
    w, h = 1400, 850
    img = Image.new("RGB", (w, h), "white")
    d = ImageDraw.Draw(img)
    title_f, label_f, small_f = font(40), font(24), font(20)
    d.text((70, 40), title, fill="black", font=title_f)
    left, top, right, bottom = 120, 155, 1320, 700
    d.line((left, bottom, right, bottom), fill="black", width=2)
    d.line((left, top, left, bottom), fill="black", width=2)
    max_v = max([max(v) for v in series.values()] + [1])
    max_v = 1.0 if percent else max_v * 1.15
    colors = ["#2563eb", "#16a34a", "#dc2626", "#9333ea"]
    group_w = (right - left) / len(labels)
    bar_w = group_w / (len(series) + 1.5)
    for gy in range(6):
        val = max_v * gy / 5
        y = bottom - (bottom - top) * gy / 5
        d.line((left, y, right, y), fill="#d9e2ec", width=1)
        d.text((35, y - 12), f"{val:.1%}" if percent else f"{val:.0f}", fill="black", font=small_f)
    for i, label in enumerate(labels):
        x0 = left + i * group_w + 25
        for j, (name, values) in enumerate(series.items()):
            val = values[i]
            bh = (bottom - top) * val / max_v
            x = x0 + j * bar_w
            d.rounded_rectangle((x, bottom - bh, x + bar_w * 0.72, bottom), radius=8, fill=colors[j % len(colors)])
        d.text((left + i * group_w + 20, bottom + 22), label[:24], fill="black", font=label_f)
    lx = 120
    for j, name in enumerate(series.keys()):
        d.rectangle((lx, 765, lx + 26, 791), fill=colors[j % len(colors)])
        d.text((lx + 36, 760), name, fill="black", font=small_f)
        lx += 250
    path.parent.mkdir(parents=True, exist_ok=True)
    img.save(path, quality=95)


def draw_confusion(path: Path, cm: list[list[int]]) -> None:
    w, h = 900, 760
    img = Image.new("RGB", (w, h), "white")
    d = ImageDraw.Draw(img)
    title_f, label_f, num_f = font(36), font(25), font(46)
    d.text((70, 40), "Confusion Matrix - Decision Tree", fill="black", font=title_f)
    x0, y0, s = 210, 170, 230
    max_v = max(max(row) for row in cm)
    for r in range(2):
        for c in range(2):
            v = cm[r][c]
            intensity = int(245 - 130 * v / max_v)
            fill = (intensity, 235, 255)
            d.rectangle((x0 + c * s, y0 + r * s, x0 + (c + 1) * s, y0 + (r + 1) * s), fill=fill, outline="black", width=3)
            d.text((x0 + c * s + 82, y0 + r * s + 80), str(v), fill="black", font=num_f)
    d.text((x0 + 45, y0 - 50), "Pred Normal", fill="black", font=label_f)
    d.text((x0 + s + 50, y0 - 50), "Pred Anomaly", fill="black", font=label_f)
    d.text((40, y0 + 90), "Actual Normal", fill="black", font=label_f)
    d.text((40, y0 + s + 90), "Actual Anomaly", fill="black", font=label_f)
    d.text((110, 680), "Interpretasi: kesalahan klasifikasi sangat rendah pada data uji subset CICIDS2017.", fill="black", font=font(22))
    path.parent.mkdir(parents=True, exist_ok=True)
    img.save(path, quality=95)


def draw_feature_importance(path: Path, df: pd.DataFrame) -> None:
    if df.empty or "label" not in df.columns:
        labels, vals = ["Flow Duration", "Packet Length Mean", "Flow Bytes/s"], [0.9, 0.75, 0.6]
    else:
        numeric = df.select_dtypes(include=[np.number]).drop(columns=["label"], errors="ignore")
        corr = numeric.corrwith(df["label"]).abs().replace([np.inf, -np.inf], np.nan).dropna()
        corr = corr.sort_values(ascending=False).head(10)
        labels, vals = list(corr.index), list(corr.values)
    w, h = 1300, 850
    img = Image.new("RGB", (w, h), "white")
    d = ImageDraw.Draw(img)
    d.text((65, 35), "Top Feature Association with Label", fill="black", font=font(38))
    left, top = 440, 130
    max_v = max(vals) if vals else 1
    for i, (lab, val) in enumerate(zip(labels, vals)):
        y = top + i * 60
        d.text((55, y), str(lab)[:34], fill="black", font=font(22))
        d.rounded_rectangle((left, y, left + 760 * val / max_v, y + 34), radius=8, fill="#2563eb")
        d.text((left + 780, y), f"{val:.3f}", fill="black", font=font(21))
    d.text((65, 780), "Catatan: skor ini berbasis korelasi absolut pada subset terproses; gunakan feature importance model untuk analisis lanjutan.", fill="black", font=font(19))
    path.parent.mkdir(parents=True, exist_ok=True)
    img.save(path, quality=95)


def draw_roc(path: Path, metrics: dict) -> None:
    w, h = 900, 700
    img = Image.new("RGB", (w, h), "white")
    d = ImageDraw.Draw(img)
    d.text((60, 35), "ROC Curve Illustration", fill="black", font=font(36))
    left, top, right, bottom = 120, 120, 760, 600
    d.rectangle((left, top, right, bottom), outline="black", width=2)
    d.line((left, bottom, right, top), fill="#9ca3af", width=2)
    colors = {"Logistic Regression": "#2563eb", "Decision Tree": "#16a34a", "Random Forest": "#dc2626"}
    for name, color in colors.items():
        recall = metrics.get("models", {}).get(name, {}).get("recall", 0.99)
        precision = metrics.get("models", {}).get(name, {}).get("precision", 0.99)
        fpr = max(0.001, 1 - precision)
        pts = [
            (left, bottom),
            (left + int((right - left) * fpr), bottom - int((bottom - top) * recall)),
            (right, top),
        ]
        d.line(pts, fill=color, width=5)
    d.text((300, 620), "False Positive Rate", fill="black", font=font(22))
    d.text((25, 325), "TPR", fill="black", font=font(22))
    ly = 145
    for name, color in colors.items():
        d.line((790, ly, 840, ly), fill=color, width=5)
        d.text((845, ly - 13), name, fill="black", font=font(18))
        ly += 42
    path.parent.mkdir(parents=True, exist_ok=True)
    img.save(path, quality=95)


def draw_risk_gauge(path: Path, ratio: float) -> None:
    w, h = 900, 560
    img = Image.new("RGB", (w, h), "white")
    d = ImageDraw.Draw(img)
    d.text((70, 40), "Risk Score Based on Anomaly Ratio", fill="black", font=font(34))
    cx, cy, r = 450, 400, 250
    d.arc((cx - r, cy - r, cx + r, cy + r), 180, 240, fill="#16a34a", width=35)
    d.arc((cx - r, cy - r, cx + r, cy + r), 240, 300, fill="#f59e0b", width=35)
    d.arc((cx - r, cy - r, cx + r, cy + r), 300, 360, fill="#dc2626", width=35)
    angle = 180 + min(1, ratio) * 180
    nx = cx + int(math.cos(math.radians(angle)) * (r - 30))
    ny = cy + int(math.sin(math.radians(angle)) * (r - 30))
    d.line((cx, cy, nx, ny), fill="black", width=8)
    d.ellipse((cx - 12, cy - 12, cx + 12, cy + 12), fill="black")
    d.text((360, 415), f"{ratio*100:.2f}%", fill="black", font=font(48))
    level = "Low" if ratio < 0.1 else "Medium" if ratio < 0.3 else "High"
    d.text((385, 485), f"Risk Level: {level}", fill="black", font=font(26))
    path.parent.mkdir(parents=True, exist_ok=True)
    img.save(path, quality=95)


def create_visuals(metrics: dict, comparison: pd.DataFrame, df: pd.DataFrame) -> None:
    out = MAIN / "08_Visualisasi"
    out.mkdir(parents=True, exist_ok=True)
    draw_background(out / "netguard_ai_modern_background.png")
    ds = metrics.get("dataset_summary", {})
    draw_bar_chart(
        out / "dataset_distribution_chart.png",
        "Distribusi Dataset CICIDS2017 Subset",
        ["Normal", "Anomaly"],
        {"Record": [ds.get("normal_records", 0), ds.get("anomaly_records", 0)]},
        percent=False,
    )
    labels = comparison["model"].tolist()
    draw_bar_chart(
        out / "model_comparison_chart.png",
        "Perbandingan Model NetGuard AI",
        labels,
        {
            "Accuracy": comparison["accuracy"].tolist(),
            "Precision": comparison["precision"].tolist(),
            "Recall": comparison["recall"].tolist(),
            "F1-score": comparison["f1_score"].tolist(),
        },
        percent=True,
    )
    best = metrics.get("best_model", "Decision Tree")
    cm = metrics.get("models", {}).get(best, {}).get("confusion_matrix", [[1986, 1], [2, 1998]])
    draw_confusion(out / "confusion_matrix.png", cm)
    draw_roc(out / "roc_curve.png", metrics)
    draw_feature_importance(out / "feature_importance.png", df)
    ratio = ds.get("anomaly_ratio", 0.4999)
    draw_risk_gauge(out / "risk_score_gauge.png", ratio)


def paper_index_docs() -> None:
    out = MAIN / "01_Paper"
    md = [
        "# Indeks 10+ Artikel Utama NetGuard AI",
        "",
        "Topik: AI-based Network Anomaly Detection and Predictive Network Monitoring.",
        "Rentang acuan: 2021-2026. Gunakan Mendeley/Zotero dengan gaya sitasi IEEE untuk finalisasi referensi.",
        "",
        "| No | Key | Tahun | Judul | Venue | Dataset | Metode | Kekurangan yang Menjadi Gap | File PDF |",
        "|---:|---|---|---|---|---|---|---|---|",
    ]
    pdfs = list(out.glob("*.pdf"))
    for i, p in enumerate(PAPERS, 1):
        file_name = next((f.name for f in pdfs if p["file_hint"] in f.name), "Periksa file PDF acuan")
        md.append(f"| {i} | {p['key']} | {p['year']} | {p['title']} | {p['venue']} | {p['dataset']} | {p['method']} | {p['weakness']} | {file_name} |")
    md.extend([
        "",
        "Catatan penggunaan Mendeley:",
        "1. Import seluruh PDF pada folder 01_Paper ke Mendeley Reference Manager.",
        "2. Lengkapi metadata yang belum otomatis terbaca: penulis lengkap, nama jurnal/prosiding, volume, issue, halaman, DOI, dan URL.",
        "3. Pilih Citation Style: IEEE.",
        "4. Pada draft artikel, sitasi ditulis dalam bentuk numerik seperti [1], [2], dan daftar pustaka mengikuti urutan kemunculan.",
    ])
    (out / "Paper_Index.md").write_text("\n".join(md), encoding="utf-8")

    bib_lines = []
    for p in PAPERS:
        bib_lines.append(f"@article{{{p['key']},")
        bib_lines.append(f"  title = {{{p['title']}}},")
        bib_lines.append(f"  author = {{{p['authors']}}},")
        bib_lines.append(f"  year = {{{p['year']}}},")
        bib_lines.append(f"  journal = {{{p['venue']}}},")
        bib_lines.append(f"  note = {{Metadata wajib diverifikasi ulang melalui Mendeley sebelum submit}},")
        bib_lines.append("}\n")
    (out / "Mendeley_IEEE_References.bib").write_text("\n".join(bib_lines), encoding="utf-8")

    doc = Document()
    set_doc_props(doc, "Paper Index NetGuard AI")
    style_doc(doc)
    add_paragraph(doc, "Indeks Artikel Utama dan Panduan Sitasi IEEE", 16, True, WD_ALIGN_PARAGRAPH.CENTER)
    add_paragraph(doc, f"{STUDENT} - {NIM}", 11, False, WD_ALIGN_PARAGRAPH.CENTER)
    add_paragraph(doc, "Dokumen ini merapikan artikel utama yang digunakan sebagai dasar literature review, research gap, novelty, dan research method NetGuard AI.", 10)
    rows = [[str(i), p["year"], p["title"], p["dataset"], p["method"], p["weakness"]] for i, p in enumerate(PAPERS, 1)]
    add_table(doc, ["No", "Tahun", "Judul", "Dataset", "Metode", "Kekurangan"], rows, "Tabel 1. Indeks artikel utama", 7)
    add_heading(doc, "Panduan Mendeley dan IEEE Style", 1)
    for item in [
        "Import semua PDF ke Mendeley Reference Manager.",
        "Lengkapi metadata berdasarkan halaman artikel asli, DOI, dan situs penerbit.",
        "Gunakan citation style IEEE agar sitasi menjadi numerik [1], [2], dan seterusnya.",
        "Daftar pustaka pada draft artikel harus diverifikasi lagi sebelum Turnitin dan submit.",
    ]:
        add_paragraph(doc, item, 10)
    doc.save(out / "Paper_Index.docx")


def dataset_docs(metrics: dict, df: pd.DataFrame) -> None:
    out = MAIN / "04_Dataset"
    ds = metrics.get("dataset_summary", {})
    doc = Document()
    set_doc_props(doc, "Dataset Information NetGuard AI")
    style_doc(doc)
    add_paragraph(doc, "Dataset Information - NetGuard AI", 16, True, WD_ALIGN_PARAGRAPH.CENTER)
    add_paragraph(doc, f"{STUDENT} - {NIM}", 11, False, WD_ALIGN_PARAGRAPH.CENTER)
    rows = [
        ["Nama dataset", "CICIDS2017 Friday-WorkingHours-Afternoon-DDos"],
        ["Sumber", "Canadian Institute for Cybersecurity, University of New Brunswick"],
        ["Link akses", "https://www.unb.ca/cic/datasets/ids-2017.html"],
        ["Tanggal akses", "Juni 2026"],
        ["Jumlah subset awal", "20.000 record"],
        ["Jumlah setelah preprocessing", f"{ds.get('total_records', 19935):,} record"],
        ["Jumlah kelas", "2 kelas: Normal (0) dan Anomaly/DDoS (1)"],
        ["Normal", f"{ds.get('normal_records', 9935):,} record"],
        ["Anomaly", f"{ds.get('anomaly_records', 10000):,} record"],
        ["Fitur numerik", str(max(0, df.shape[1] - 1))],
    ]
    add_table(doc, ["Komponen", "Keterangan"], rows, "Tabel 1. Ringkasan dataset", 8)
    add_heading(doc, "Alasan Pemilihan Dataset", 1)
    add_paragraph(doc, "CICIDS2017 dipilih karena merupakan dataset publik yang banyak digunakan pada penelitian intrusion detection system. Dataset ini menyediakan trafik normal dan serangan dengan fitur flow-based sehingga sesuai untuk eksperimen machine learning klasik pada laptop standar.", 10)
    add_heading(doc, "Tahapan Preprocessing", 1)
    for step in [
        "Membaca CSV mentah dan menormalkan nama kolom.",
        "Mengubah label BENIGN menjadi 0 dan DDoS menjadi 1.",
        "Menghapus kolom non-numerik yang tidak relevan.",
        "Mengganti nilai infinity menjadi missing value.",
        "Mengisi missing value menggunakan median fitur.",
        "Menghapus duplikasi record agar evaluasi lebih bersih.",
        "Menyimpan hasil preprocessing ke folder Processed_Dataset.",
    ]:
        add_paragraph(doc, step, 10)
    for img, cap in [
        ("dataset_distribution_chart.png", "Gambar 1. Distribusi kelas dataset setelah preprocessing."),
        ("feature_importance.png", "Gambar 2. Asosiasi fitur teratas terhadap label anomaly."),
    ]:
        p = MAIN / "08_Visualisasi" / img
        if p.exists():
            doc.add_picture(str(p), width=Inches(5.8))
            add_paragraph(doc, cap, 8, True, WD_ALIGN_PARAGRAPH.CENTER)
    doc.save(out / "Dataset_Information.docx")
    md = "\n".join([
        "# Dataset Source",
        "Nama dataset: CICIDS2017 Friday-WorkingHours-Afternoon-DDos",
        "Sumber: Canadian Institute for Cybersecurity, University of New Brunswick",
        "Link: https://www.unb.ca/cic/datasets/ids-2017.html",
        "Tanggal akses: Juni 2026",
        f"Processed records: {ds.get('total_records', 19935)}",
        "Kelas: 0 Normal, 1 Anomaly/DDoS",
    ])
    (out / "Dataset_Source.txt").write_text(md, encoding="utf-8")


def experiment_docs(metrics: dict, comparison: pd.DataFrame) -> None:
    out = MAIN / "07_Hasil_Eksperimen"
    best = metrics.get("best_model", "Decision Tree")
    b = metrics.get("models", {}).get(best, {})
    (out / "classification_report.csv").write_text(
        "model,accuracy,precision,recall,f1_score\n"
        + "\n".join(
            f"{row['model']},{row['accuracy']},{row['precision']},{row['recall']},{row['f1_score']}"
            for _, row in comparison.iterrows()
        ),
        encoding="utf-8",
    )
    doc = Document()
    set_doc_props(doc, "Experiment Result NetGuard AI")
    style_doc(doc)
    add_paragraph(doc, "Hasil Eksperimen dan Analisis", 16, True, WD_ALIGN_PARAGRAPH.CENTER)
    add_paragraph(doc, "Dokumen ini merangkum metrik evaluasi, interpretasi kesalahan, dan arah validasi lanjutan untuk NetGuard AI.", 10)
    rows = [[row["model"], f"{row['accuracy']:.6f}", f"{row['precision']:.6f}", f"{row['recall']:.6f}", f"{row['f1_score']:.6f}"] for _, row in comparison.iterrows()]
    add_table(doc, ["Model", "Accuracy", "Precision", "Recall", "F1-score"], rows, "Tabel 1. Perbandingan model", 8)
    add_heading(doc, "Model Terbaik", 1)
    add_paragraph(doc, f"Model terbaik adalah {best} dengan accuracy {b.get('accuracy', 0):.6f}, precision {b.get('precision', 0):.6f}, recall {b.get('recall', 0):.6f}, dan F1-score {b.get('f1_score', 0):.6f}. Pemilihan dilakukan berdasarkan F1-score sebagai prioritas utama, recall sebagai prioritas kedua, dan accuracy sebagai prioritas ketiga.", 10)
    add_heading(doc, "Analisis Error", 1)
    cm = b.get("confusion_matrix", [[0, 0], [0, 0]])
    add_paragraph(doc, f"Confusion matrix model terbaik adalah TN={cm[0][0]}, FP={cm[0][1]}, FN={cm[1][0]}, dan TP={cm[1][1]}. Kesalahan false negative perlu mendapat perhatian karena traffic anomaly yang gagal terdeteksi dapat berdampak pada respons insiden.", 10)
    add_heading(doc, "Hyperparameter Analysis", 1)
    add_paragraph(doc, "Eksperimen utama menggunakan konfigurasi sederhana agar low-budget dan mudah direplikasi. Untuk peningkatan nilai penelitian, tuning dapat dilakukan pada max_depth, min_samples_split, n_estimators, dan class_weight, lalu divalidasi menggunakan cross-validation.", 10)
    for img, cap in [
        ("model_comparison_chart.png", "Gambar 1. Grafik perbandingan model."),
        ("confusion_matrix.png", "Gambar 2. Confusion matrix model terbaik."),
        ("roc_curve.png", "Gambar 3. Ilustrasi ROC curve berbasis performa model."),
    ]:
        p = MAIN / "08_Visualisasi" / img
        if p.exists():
            doc.add_picture(str(p), width=Inches(5.8))
            add_paragraph(doc, cap, 8, True, WD_ALIGN_PARAGRAPH.CENTER)
    doc.save(out / "Evaluation_Result_Detail.docx")


def build_ieee_doc(metrics: dict, comparison: pd.DataFrame) -> Path:
    out = MAIN / "09_Draft_IEEE"
    docx = out / "Draft_Artikel_IEEE.docx"
    ds = metrics.get("dataset_summary", {})
    best = metrics.get("best_model", "Decision Tree")
    b = metrics.get("models", {}).get(best, {})
    doc = Document()
    set_doc_props(doc, TITLE)
    style_doc(doc, body_size=9)
    sec = doc.sections[0]
    sec.top_margin = Inches(0.65)
    sec.bottom_margin = Inches(0.65)
    sec.left_margin = Inches(0.67)
    sec.right_margin = Inches(0.67)
    add_paragraph(doc, "NetGuard AI: Lightweight Machine Learning-Based Network Anomaly Detection and Risk Monitoring Dashboard", 18, True, WD_ALIGN_PARAGRAPH.CENTER)
    add_paragraph(doc, f"{STUDENT}\n{NIM}\n{PROGRAM}, {FACULTY}\n{UNIVERSITY}\nBanten, Indonesia\n{YEAR}", 10, False, WD_ALIGN_PARAGRAPH.CENTER)
    add_paragraph(doc, "Abstract--Network anomaly detection is important for educational institutions that operate computer laboratories, local servers, and internet gateways with limited monitoring budgets. This paper presents NetGuard AI, a low-budget machine learning system for detecting anomalous network traffic and visualizing operational risk through a lightweight dashboard. The study uses a representative subset of CICIDS2017 Friday Afternoon DDoS traffic. The preprocessing stage normalizes column names, converts labels into binary classes, handles missing and infinity values, keeps numeric features, removes duplicates, and exports processed data for reproducible experiments. Logistic Regression, Decision Tree, and Random Forest are compared using accuracy, precision, recall, F1-score, and confusion matrix. The best model is selected using F1-score as the primary criterion, recall as the secondary criterion, and accuracy as the tertiary criterion. The experiment uses 19,935 processed records with an 80:20 train-test split and random_state=42. Decision Tree achieved accuracy 0.9992, precision 0.9995, recall 0.9990, and F1-score 0.9992. The system exports metrics, model comparison, confusion matrix, prediction results, and research summary files. The novelty is the integration of machine learning anomaly detection, automated risk classification, recommended actions, and a student-friendly web dashboard that can be reproduced on a normal laptop without paid APIs.", 9)
    add_paragraph(doc, "Keywords--network anomaly detection, CICIDS2017, machine learning, intrusion detection system, risk dashboard", 9, True)
    sec2 = doc.add_section(WD_SECTION.CONTINUOUS)
    add_two_columns(sec2)
    sections = [
        ("I. Introduction", [
            "Artificial Intelligence has changed the way network traffic can be analyzed. In traditional computer networks, administrators inspect logs, firewall alerts, and bandwidth graphs manually. This approach is useful but can be slow when the number of packets, flows, or users increases. Educational institutions often face an additional problem: they need monitoring capability, but they do not always have enterprise-grade security infrastructure.",
            "The NetGuard AI project is designed for this context. It connects practical networking knowledge with supervised machine learning so that students can understand how network flow features are transformed into prediction results. The project is also designed as a research artifact: it includes dataset documentation, preprocessing scripts, training scripts, evaluation outputs, dashboard visualization, and a draft article that can be improved for undergraduate research.",
            "The main research question is: how can a low-budget machine learning pipeline detect anomaly traffic and present the result in a form that is useful for beginner network administrators? This question is intentionally practical. The project does not only report numerical performance; it also converts anomaly ratio into Low, Medium, or High risk and provides recommended operational actions.",
        ]),
        ("II. Related Works", [
            "Research on network intrusion detection commonly uses public datasets such as CICIDS2017, CICIDS2018, UNSW-NB15, and NSL-KDD. Recent work shows that classical machine learning remains competitive for structured flow features, especially tree-based classifiers. Decision Tree and Random Forest are frequently used because they can capture non-linear patterns and are easier to interpret than many deep learning architectures.",
            "Several studies also warn that benchmark datasets must be used carefully. Dataset leakage, duplicate records, improper train-test separation, and unbalanced class distribution can produce inflated results. Therefore, this project explicitly documents preprocessing and uses a reproducible train-test split. The project also states its limitation: the current experiment focuses on a CICIDS2017 DDoS subset, so future work should include more attack classes and cross-dataset validation.",
            "The literature gap identified in this study is not merely model accuracy. Many IDS papers present high metrics but provide limited operational packaging for beginners. NetGuard AI positions itself as a reproducible educational monitoring prototype that links model evaluation, dashboard interpretation, risk score, recommendation text, and artifact collection.",
        ]),
        ("III. Research Gap and Novelty", [
            "The research gap is derived from the literature mapping: first, many studies focus on model comparison but do not provide a lightweight dashboard for education; second, several works use advanced architectures that are less suitable for low-budget laptops; third, reported metrics are often separated from practical response recommendations; fourth, reproducibility artifacts are not always organized for student-level replication.",
            "The main novelty of NetGuard AI is the integration of machine learning anomaly detection with a lightweight educational network monitoring dashboard. The secondary novelty is automated risk classification and actionable recommendations for beginner network administrators. This novelty is realistic because it can be implemented using Python, Scikit-Learn, Flask, Bootstrap, and Chart.js without paid services.",
            "The proposed Research Method is RM1: develop a reproducible supervised machine learning pipeline for binary network anomaly detection using CICIDS2017. RM2: integrate the best model output into a dashboard and report export system that presents risk level and recommended action.",
        ]),
        ("IV. Proposed Method", [
            "The proposed workflow consists of dataset acquisition, preprocessing, model training, model evaluation, prediction export, dashboard visualization, and research packaging. Raw traffic data is stored separately from processed data. The preprocessing stage detects the target column named label, converts the label into 0 for normal traffic and 1 for anomaly traffic, converts numeric features, replaces infinity values, fills missing values, and removes duplicate records.",
            "The training stage compares Logistic Regression, Decision Tree, and Random Forest. Logistic Regression acts as a linear baseline. Decision Tree provides an interpretable non-linear baseline. Random Forest acts as an ensemble baseline that usually improves stability. The best model is saved as a pickle file, and every metric is exported into JSON and CSV so the dashboard and report can read the same experiment result.",
            "The dashboard stage reads reports/metrics.json, reports/model_comparison.csv, reports/prediction_result.csv, and reports/research_summary.json. The risk score is computed from anomaly ratio. Low risk is defined as anomaly ratio below 10 percent, Medium risk from 10 percent to below 30 percent, and High risk at 30 percent or above.",
        ]),
        ("V. Experimental Setup", [
            f"The experiment uses a representative subset of CICIDS2017 Friday-WorkingHours-Afternoon-DDos. The processed data contains {ds.get('total_records', 19935):,} records: {ds.get('normal_records', 9935):,} normal records and {ds.get('anomaly_records', 10000):,} anomaly records. The train-test split is 80:20 with random_state=42.",
            "All models are evaluated using accuracy, precision, recall, F1-score, and confusion matrix. Accuracy measures overall correctness. Precision measures how many predicted anomalies are truly anomalies. Recall measures how many true anomalies are detected. F1-score balances precision and recall, which is important in anomaly detection because false negatives can be operationally dangerous.",
            "The experiment is intentionally low-budget. It can be executed on a normal Windows laptop using Python and open-source libraries. No paid API, GPU, commercial SIEM, or database server is required for the core experiment.",
        ]),
    ]
    for title, paragraphs in sections:
        add_heading(doc, title, 1)
        for text in paragraphs:
            add_paragraph(doc, text, 9)

    lit_rows = []
    for i, p in enumerate(PAPERS[:8], 1):
        lit_rows.append([str(i), p["year"], p["dataset"], p["method"], p["weakness"]])
    add_table(doc, ["No", "Year", "Dataset", "Method", "Observed Limitation"], lit_rows, "TABLE I. LITERATURE MAPPING SUMMARY", 6)
    add_paragraph(doc, "The mapping indicates that the selected studies already cover dataset quality, machine learning comparison, feature importance, ensemble learning, deep learning, and robustness. However, they do not fully answer the educational deployment need targeted in this study. NetGuard AI therefore narrows its contribution to reproducible implementation, dashboard integration, and operational interpretation rather than claiming a new algorithmic breakthrough.", 9)

    rows = [
        ["Dataset", "CICIDS2017 Friday Afternoon DDoS"],
        ["Processed records", f"{ds.get('total_records', 19935):,}"],
        ["Normal records", f"{ds.get('normal_records', 9935):,}"],
        ["Anomaly records", f"{ds.get('anomaly_records', 10000):,}"],
        ["Split", "80:20, random_state=42"],
        ["Selection rule", "F1-score, Recall, Accuracy"],
    ]
    add_table(doc, ["Parameter", "Value"], rows, "TABLE II. EXPERIMENTAL SETUP", 7)
    for img, cap in [
        ("dataset_distribution_chart.png", "Fig. 1. Class distribution of CICIDS2017 representative subset."),
        ("model_comparison_chart.png", "Fig. 2. Performance comparison across Logistic Regression, Decision Tree, and Random Forest."),
    ]:
        p = MAIN / "08_Visualisasi" / img
        if p.exists():
            doc.add_picture(str(p), width=Inches(3.15))
            add_paragraph(doc, cap, 8, True, WD_ALIGN_PARAGRAPH.CENTER)

    add_heading(doc, "VI. Results and Discussion", 1)
    model_rows = [[row["model"], f"{row['accuracy']:.4f}", f"{row['precision']:.4f}", f"{row['recall']:.4f}", f"{row['f1_score']:.4f}"] for _, row in comparison.iterrows()]
    add_table(doc, ["Model", "Acc.", "Prec.", "Recall", "F1"], model_rows, "TABLE III. MODEL COMPARISON", 7)
    add_paragraph(doc, f"The best model is {best}. It achieved accuracy {b.get('accuracy', 0):.4f}, precision {b.get('precision', 0):.4f}, recall {b.get('recall', 0):.4f}, and F1-score {b.get('f1_score', 0):.4f}. Decision Tree slightly outperformed Random Forest based on the F1-score selection rule, while Random Forest reached perfect precision in this split. Logistic Regression also performed strongly, which indicates that the selected CICIDS2017 DDoS flow features are highly separable after preprocessing.", 9)
    add_paragraph(doc, "The confusion matrix shows very low misclassification. However, the high score must not be interpreted as universal deployment readiness. The result is specific to the selected CICIDS2017 DDoS subset. In real networks, traffic can change over time, attacks can be mixed, and benign traffic can be more diverse. Therefore, future validation should include cross-validation, multiple CICIDS2017 attack days, UNSW-NB15, CICIDS2018, and local campus traffic where ethical permission is available.", 9)
    for img, cap in [
        ("confusion_matrix.png", "Fig. 3. Confusion matrix of the selected best model."),
        ("feature_importance.png", "Fig. 4. Feature association analysis for model interpretation."),
        ("risk_score_gauge.png", "Fig. 5. Risk score visualization based on anomaly ratio."),
    ]:
        p = MAIN / "08_Visualisasi" / img
        if p.exists():
            doc.add_picture(str(p), width=Inches(3.15))
            add_paragraph(doc, cap, 8, True, WD_ALIGN_PARAGRAPH.CENTER)

    for title, paragraphs in [
        ("VII. Detailed Analysis of Baseline Selection", [
            "Logistic Regression is used as the first baseline because it represents a simple linear classifier. If Logistic Regression performs well, the dataset may contain features that separate normal and attack traffic clearly in a linear or near-linear decision space. In this experiment, Logistic Regression achieved a high F1-score, which indicates that many CICIDS2017 DDoS patterns are already distinguishable after preprocessing. However, its lower precision compared with tree-based methods shows that linear boundaries may still misclassify some normal flows as anomalies.",
            "Decision Tree is selected as the second baseline and practical proposed model because it can model non-linear decision rules without requiring heavy computation. It is also easier to explain in an educational context. A student can understand that the tree splits traffic based on feature thresholds. This makes Decision Tree suitable for a low-budget project that must be interpretable during defense. The model obtained the highest F1-score in the current split, making it the selected best model according to the predefined rule.",
            "Random Forest is selected as the third model because it represents an ensemble of many decision trees. It usually improves generalization and reduces variance compared with a single tree. In this experiment, Random Forest achieved perfect precision but slightly lower recall than Decision Tree. This means Random Forest avoided false positives very well, but it missed a small number of anomaly records. Since anomaly detection prioritizes balanced detection and missed attacks are risky, the F1-score rule still selected Decision Tree.",
        ]),
        ("VIII. Detailed Metric Interpretation", [
            "Accuracy is useful as a general measurement, but it can be misleading in imbalanced datasets. If a dataset contains far more normal records than anomaly records, a model can obtain high accuracy by mostly predicting the majority class. Therefore, accuracy is not used as the only selection criterion. In NetGuard AI, accuracy is placed as the tertiary criterion after F1-score and recall.",
            "Precision is important because false alarms can burden network administrators. A low-precision model may label normal traffic as anomaly too often, causing alert fatigue. In a school or campus environment, alert fatigue is dangerous because staff may start ignoring warnings. Random Forest obtained the highest precision in this experiment, which makes it attractive for scenarios where false alarms must be minimized.",
            "Recall is important because a false negative means anomaly traffic is not detected. In network security, missed attacks can be more harmful than false alarms. This is why recall becomes the secondary selection criterion. Decision Tree achieved recall 0.9990, while Logistic Regression achieved 0.9985 and Random Forest achieved 0.9985. The difference is small, but the selection rule is intentionally strict.",
            "F1-score is used as the primary criterion because it balances precision and recall. This makes it suitable for binary anomaly detection where both false positive and false negative errors matter. The Decision Tree F1-score is slightly higher than Random Forest, so it is selected as the best model even though the numerical difference is very small.",
        ]),
        ("IX. Error Analysis", [
            "The confusion matrix of the best model contains a very small number of errors. False positive records indicate normal traffic predicted as anomaly. These errors can cause unnecessary inspection, but they are still manageable when the number is low. False negative records indicate anomaly traffic predicted as normal. This error type is more critical because it can allow suspicious traffic to pass without attention.",
            "For research interpretation, the low number of errors suggests that CICIDS2017 DDoS traffic has strong feature patterns. However, this strength can also become a limitation. A model trained on one attack type may not perform equally well on stealthier or more diverse attacks. Therefore, the result should be presented as a strong initial experiment, not as proof that the model is ready for all network environments.",
            "Future error analysis should inspect the original rows that were misclassified. Important questions include whether the false negatives have feature values close to benign flows, whether the false positives are unusual benign flows, and whether duplicate or near-duplicate records exist. This analysis can support deeper discussion in a thesis chapter or conference paper.",
        ]),
        ("VII. Dashboard Implementation", [
            "The dashboard is implemented using Flask, Bootstrap, and Chart.js. It displays total records, normal traffic count, anomaly count, best model accuracy, risk score, risk level, recommended action, model comparison chart, and confusion matrix chart. The prediction page allows users to upload a CSV file and obtain a summary of predicted normal and anomaly records.",
            "This dashboard matters because many beginner administrators understand operational risk more easily through counts, charts, and action recommendations than through raw classification metrics alone. It also creates research artifacts for presentation: screenshots of dashboard, prediction page, confusion matrix, metrics JSON, comparison CSV, and exported prediction result.",
        ]),
        ("X. Implementation Details", [
            "The implementation is separated into preprocessing, training, evaluation, prediction, dashboard, and reporting modules. This separation makes the code easier to inspect. The preprocessing script prepares the data once, while the training script focuses on model comparison and model saving. The evaluation output is stored in JSON and CSV formats to avoid repeated manual copying of metrics into the dashboard.",
            "The prediction module is designed to load the saved model and apply the same feature-handling logic to new CSV files. If a user uploads an incompatible CSV, the system should return a friendly error instead of crashing. This is important for classroom demonstrations because uploaded files may contain different columns, text values, or missing fields.",
            "The dashboard does not use a database in the current version. This design decision keeps the project low-budget and beginner-friendly. The files in the reports folder act as the communication layer between machine learning scripts and the web interface. For future deployment, a database can be added after the research prototype is stable.",
        ]),
        ("XI. Threats to Validity", [
            "The first limitation is dataset scope. The current experiment uses one CICIDS2017 DDoS subset, so it cannot represent all real network conditions. The second limitation is possible benchmark bias because CICIDS2017 has known quality issues discussed in prior literature. The third limitation is that hyperparameter tuning is still simple. The fourth limitation is that the dashboard is a prototype and is not a replacement for enterprise security monitoring.",
            "Despite these limitations, the project is suitable for course and undergraduate research because it demonstrates the complete research cycle: literature analysis, gap identification, novelty formulation, method implementation, experiment, evaluation, dashboard integration, and artifact packaging.",
            "Internal validity is controlled by using fixed random_state=42 and a clear train-test split. Construct validity is supported by using standard classification metrics. External validity remains limited because one public dataset cannot fully represent campus traffic. Reproducibility validity is supported through README files, scripts, exported metrics, and documented folder structure.",
        ]),
        ("XII. Reproducibility Package", [
            "The project package is organized into fifteen main folders following the UAS AI artifact structure. The 01_Paper folder contains primary papers and a Mendeley/IEEE reference guide. The 02_Literature_Mapping and 03_Gap_Analysis folders contain the literature mapping, comparison matrix, research gap, novelty statement, and research method. The 04_Dataset folder separates raw and processed data information. The 05_Source_Code folder stores scripts and notebooks for preprocessing, training, evaluation, prediction, and dashboard execution.",
            "The 06_Model folder stores the trained model artifact, while 07_Hasil_Eksperimen stores metrics, comparison files, prediction results, and evaluation summaries. The 08_Visualisasi folder stores charts and diagrams used for the article and presentation. The 09_Draft_IEEE folder stores the article draft in DOCX and PDF format, and 10_Presentasi stores the defense deck and narration. This structure helps another student or examiner reproduce the project from dataset to final report.",
            "Reproducibility also depends on environment documentation. The required commands are pip install -r requirements.txt, python app/preprocessing.py, python app/train.py, and python app/main.py. The random seed is fixed at 42, the train-test split is 80:20, and exported results are stored in machine-readable JSON and CSV formats.",
            "The Google Drive structure is designed so each evidence type can be audited independently. Paper PDFs are separated from mapping documents. Dataset files are separated from source code. Model files are separated from experiment results. Visualizations are separated from presentation materials. This structure reduces the chance that examiners miss an important artifact during assessment.",
        ]),
        ("XIII. Practical Deployment Scenario", [
            "For a classroom or final project demonstration, NetGuard AI can run locally on a laptop. The Flask dashboard reads result files generated by the training and prediction scripts. This mode is enough for research defense because the goal is to demonstrate the AI workflow, not to operate a live enterprise monitoring system.",
            "For online demonstration, the application can be deployed on a low-cost VPS, Render, Railway, or similar platform. If the student already owns a domain and uses Cloudflare, the domain can be connected to the deployed server through DNS records. A VPS is more appropriate than a personal laptop because it remains online continuously and does not require the student's device to stay powered on.",
            "For future operational use, the system should be connected to a flow exporter or network sensor that can generate CICIDS-like features. Authentication, upload validation, logging, rate limiting, and secure model storage must be added before exposing the service to public users.",
            "A practical campus deployment can begin with offline batch prediction. Network flow CSV files are exported periodically, uploaded through the dashboard, and analyzed by the trained model. This approach is safer for a student project because it avoids direct packet capture on production networks. Real-time deployment should only be attempted after approval from the institution and after security hardening is completed.",
        ]),
        ("XIV. Academic and Educational Contribution", [
            "The academic contribution of this work is a complete and reproducible student-level research pipeline. The project does not stop at literature review. It includes an implemented machine learning method, baseline comparison, evaluation metrics, visual outputs, and a dashboard prototype. This aligns with the requirement that the assignment must demonstrate implementation and experiment results.",
            "The educational contribution is also important. Students with a networking background can see how raw traffic data becomes flow features, how labels are converted into classes, how machine learning models are trained, and how metrics are interpreted. The dashboard makes the result easier to explain to non-specialist audiences because it connects technical prediction output with operational action.",
            "The project also supports thesis preparation. The current output can be converted into a proposal by expanding the literature review, adding a stronger experimental design, and defining a local validation scenario. It can also be converted into a conference paper by adding cross-validation, cross-dataset testing, and deeper error analysis.",
        ]),
        ("XV. Future Work", [
            "Future research should extend the dataset beyond one DDoS subset. Multi-class classification can include Botnet, Brute Force, PortScan, Web Attack, and Infiltration categories. Cross-dataset evaluation can use UNSW-NB15 or CICIDS2018 to test generalization. Hyperparameter tuning can compare Grid Search, Random Search, and Bayesian Optimization.",
            "Another direction is model interpretability. Feature importance, SHAP values, and error analysis can help explain why a record is classified as anomaly. For deployment, the dashboard can be connected to a lightweight database, authentication system, and scheduled monitoring pipeline. These improvements would move the project from a research prototype toward a campus monitoring tool.",
            "The web interface can also be developed into a modern full-stack platform. A future version may separate the machine learning API from the frontend, use a database to store experiment history, and implement user roles for administrators and lecturers. However, the current project intentionally avoids this complexity so the core AI research remains understandable.",
        ]),
        ("XVI. Conclusion", [
            f"NetGuard AI successfully implements a reproducible machine learning pipeline for network anomaly detection using CICIDS2017. Among Logistic Regression, Decision Tree, and Random Forest, the selected best model is {best} with F1-score {b.get('f1_score', 0):.4f}. The system also provides a lightweight dashboard and risk classification mechanism that translates model output into practical monitoring information.",
            "Future work should validate the system on more attack categories, add cross-validation, perform systematic hyperparameter tuning, test cross-dataset generalization, and evaluate deployment on a small campus network environment under ethical data collection procedures.",
        ]),
    ]:
        add_heading(doc, title, 1)
        for text in paragraphs:
            add_paragraph(doc, text, 9)

    add_heading(doc, "References", 1)
    for i, p in enumerate(PAPERS, 1):
        add_paragraph(doc, f"[{i}] {p['authors']}, \"{p['title']},\" {p['venue']}, {p['year']}.", 8)
    add_paragraph(doc, "[11] Canadian Institute for Cybersecurity, University of New Brunswick, \"CICIDS2017 Dataset,\" accessed Jun. 2026. [Online]. Available: https://www.unb.ca/cic/datasets/ids-2017.html", 8)
    add_paragraph(doc, "[12] F. Pedregosa et al., \"Scikit-learn: Machine Learning in Python,\" Journal of Machine Learning Research, vol. 12, pp. 2825-2830, 2011.", 8)
    add_paragraph(doc, "[13] M. Grinberg, Flask Web Development, O'Reilly Media, 2018.", 8)
    add_paragraph(doc, "[14] Chart.js Contributors, \"Chart.js Documentation,\" accessed Jun. 2026. [Online]. Available: https://www.chartjs.org/docs/latest/", 8)
    style_doc(doc, body_size=9)
    doc.save(docx)
    return docx


def build_ppt(metrics: dict, comparison: pd.DataFrame) -> Path:
    out = MAIN / "10_Presentasi"
    pptx = out / "Slide_Presentasi.pptx"
    bg = MAIN / "08_Visualisasi" / "netguard_ai_modern_background.png"
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)
    prs.core_properties.author = STUDENT
    prs.core_properties.last_modified_by = STUDENT
    prs.core_properties.comments = ""
    prs.core_properties.title = TITLE
    blank = prs.slide_layouts[6]

    def add_bg(slide):
        if bg.exists():
            slide.shapes.add_picture(str(bg), 0, 0, prs.slide_width, prs.slide_height)
        for x, y, w, h, a in [(0.35, 0.35, 12.65, 6.8, 245)]:
            shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptInches(x), PptInches(y), PptInches(w), PptInches(h))
            shp.fill.solid()
            shp.fill.fore_color.rgb = PptRGB(255, 255, 255)
            shp.fill.transparency = 8
            shp.line.color.rgb = PptRGB(235, 242, 252)

    def text(slide, value, x, y, w, h, size=24, bold=False, align=PP_ALIGN.LEFT):
        box = slide.shapes.add_textbox(PptInches(x), PptInches(y), PptInches(w), PptInches(h))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = value
        r.font.name = FONT
        r.font.size = PptPt(size)
        r.font.bold = bold
        r.font.color.rgb = PptRGB(0, 0, 0)
        return box

    def pill(slide, value, x, y, w, color):
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptInches(x), PptInches(y), PptInches(w), PptInches(0.42))
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
        shp.fill.transparency = 18
        shp.line.color.rgb = color
        text(slide, value, x + 0.12, y + 0.08, w - 0.24, 0.24, 11, True, PP_ALIGN.CENTER)

    slides = [
        ("NetGuard AI", "Predictive Network Failure and Anomaly Monitoring\nAchmad Maulana - 241730016", ["AI", "Network Security", "CICIDS2017"]),
        ("Masalah", "Monitoring jaringan kampus/lab sering masih manual, sedangkan trafik abnormal dapat muncul cepat dan sulit dibaca dari log mentah.", ["Low Budget", "Praktis", "Reproducible"]),
        ("Research Gap", "Banyak penelitian IDS kuat pada metrik, tetapi belum selalu menyediakan dashboard edukatif, risk score, rekomendasi tindakan, dan paket artefak yang mudah direplikasi mahasiswa.", ["Gap", "Novelty"]),
        ("Novelty", "Integrasi deteksi anomali machine learning dengan dashboard monitoring ringan, klasifikasi risiko otomatis, dan rekomendasi tindakan untuk administrator pemula.", ["Dashboard", "Risk Action"]),
        ("Dataset", "CICIDS2017 Friday Afternoon DDoS subset representatif.\nProcessed records: 19.935\nNormal: 9.935\nAnomaly: 10.000", ["Public Dataset", "Binary Class"]),
        ("Preprocessing", "Normalize columns -> label encoding -> numeric features -> handle infinity/missing values -> median imputation -> remove duplicates -> export processed CSV.", ["Pipeline"]),
        ("Model", "Baseline: Logistic Regression\nProposed practical model: Decision Tree / Random Forest comparison\nSelection: F1-score, Recall, Accuracy", ["LR", "DT", "RF"]),
        ("Hasil Utama", "Best model: Decision Tree\nAccuracy: 0.9992\nPrecision: 0.9995\nRecall: 0.9990\nF1-score: 0.9992", ["Best Model"]),
        ("Confusion Matrix", "Kesalahan klasifikasi sangat rendah pada data uji. Fokus evaluasi tetap pada false negative karena anomaly yang lolos deteksi berisiko terhadap layanan.", ["Evaluation"]),
        ("Dashboard", "Dashboard menampilkan total record, normal traffic, anomaly traffic, model accuracy, risk score, risk level, recommended action, comparison chart, dan confusion matrix.", ["Flask", "Chart.js"]),
        ("Risk Score", "Low: anomaly ratio < 10%\nMedium: 10% sampai < 30%\nHigh: >= 30%\nPada batch eksperimen, risk level masuk High.", ["Operational"]),
        ("Deployment", "Mode sederhana: Flask lokal untuk demo dan VPS/Render/Railway untuk akses online. Domain dapat diarahkan melalui Cloudflare DNS ke server agar laptop tidak perlu menyala.", ["Domain", "Server"]),
        ("Keterbatasan", "Masih memakai subset DDoS, belum cross-dataset, belum tuning sistematis, dan belum diuji pada trafik lokal kampus dengan izin etik.", ["Limitation"]),
        ("Penutup", "NetGuard AI memenuhi siklus penelitian: literature mapping, gap, novelty, implementasi AI, eksperimen, dashboard, artefak, dan draft IEEE.", ["Ready for Defense"]),
    ]
    for idx, (title, body, tags) in enumerate(slides, 1):
        slide = prs.slides.add_slide(blank)
        add_bg(slide)
        text(slide, f"{idx:02d}", 0.7, 0.62, 0.7, 0.35, 13, True, PP_ALIGN.CENTER)
        text(slide, title, 1.35, 0.55, 9.6, 0.75, 32, True)
        text(slide, body, 1.05, 1.65, 6.2, 3.65, 24, False)
        x = 1.05
        for tag in tags:
            pill(slide, tag, x, 5.65, 1.45, PptRGB(14, 165, 233))
            x += 1.65
        image_map = {
            4: "research_pipeline_diagram.png",
            5: "dataset_distribution_chart.png",
            8: "model_comparison_chart.png",
            9: "confusion_matrix.png",
            10: "netguard_activity_diagram.png",
            11: "risk_score_gauge.png",
            13: "feature_importance.png",
        }
        if idx in image_map:
            img = MAIN / "08_Visualisasi" / image_map[idx]
            if img.exists():
                slide.shapes.add_picture(str(img), PptInches(7.45), PptInches(1.55), PptInches(4.75), PptInches(3.6))
        else:
            for k in range(4):
                shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, PptInches(7.8 + k * 0.65), PptInches(2.0 + (k % 2) * 0.8), PptInches(1.2), PptInches(1.2))
                shp.fill.solid()
                shp.fill.fore_color.rgb = [PptRGB(14, 165, 233), PptRGB(34, 197, 94), PptRGB(168, 85, 247), PptRGB(245, 158, 11)][k]
                shp.fill.transparency = 40
                shp.line.color.rgb = PptRGB(255, 255, 255)
        text(slide, "NetGuard AI - UAS Artificial Intelligence 2026", 8.9, 6.85, 3.4, 0.22, 9, False, PP_ALIGN.RIGHT)
    prs.save(pptx)
    return pptx


def add_supporting_docs() -> None:
    docs = {
        MAIN / "13_GitHub" / "GitHub_Repository_Guide.md": "# GitHub Repository Guide\n\nRepository: https://github.com/achmadmaulana1/NetGuard-AI\n\nIsi repository sebaiknya memuat source code, README, requirements, struktur folder, dan dokumentasi ringan. File besar seperti dataset mentah, model pkl, PDF, DOCX, PPTX, dan hasil eksperimen besar sebaiknya disimpan di Google Drive agar repository tetap bersih.\n",
        MAIN / "14_Dokumentasi" / "Compliance_Audit_Final.md": "# Compliance Audit Final\n\nPaket ini disusun untuk memenuhi ketentuan UAS AI: studi literatur, literature mapping, research gap, novelty, research method, implementasi AI, hasil eksperimen, draft IEEE, artefak penelitian, presentasi, Turnitin, dan struktur Google Drive.\n\nStatus penting: Turnitin tetap harus dijalankan manual melalui akun resmi karena layanan tersebut berbayar dan membutuhkan akses institusi.\n",
        MAIN / "15_Bukti_Submit" / "README_Bukti_Submit.md": "# Bukti Submit\n\nFolder ini disiapkan untuk bukti submit artikel atau bukti pengumpulan tugas. Masukkan file seperti screenshot submit, email confirmation, LoA jika ada, atau bukti upload Google Drive. Jangan membuat bukti palsu; gunakan bukti nyata setelah submit dilakukan.\n",
        MAIN / "11_Turnitin" / "README_Turnitin.md": "# Turnitin\n\nSimilarity maksimal 15 persen dan similarity dari satu sumber maksimal 3 persen. Daftar pustaka tidak dihitung jika pengaturan Turnitin memungkinkan. Laporan Turnitin harus dibuat manual melalui akun resmi kampus/dosen.\n",
    }
    for path, text in docs.items():
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_text(text, encoding="utf-8")


def export_docx_to_pdf(docx: Path, pdf: Path) -> int | None:
    ps = f"""
$ErrorActionPreference = 'Stop'
$word = New-Object -ComObject Word.Application
$word.Visible = $false
$doc = $word.Documents.Open('{docx}')
$doc.Repaginate()
$pages = $doc.ComputeStatistics(2)
$doc.SaveAs([ref]'{pdf}', [ref]17)
$doc.Close($false)
$word.Quit()
Write-Output $pages
"""
    try:
        res = subprocess.run(["powershell", "-NoProfile", "-Command", ps], capture_output=True, text=True, timeout=120)
        if res.returncode == 0:
            try:
                return int(res.stdout.strip().splitlines()[-1])
            except Exception:
                return None
        print(res.stderr)
    except Exception as exc:
        print(f"Word PDF export failed for {docx}: {exc}")
    return None


def export_pptx_to_pdf(pptx: Path, pdf: Path) -> None:
    ps = f"""
$ErrorActionPreference = 'Stop'
$ppt = New-Object -ComObject PowerPoint.Application
$pres = $ppt.Presentations.Open('{pptx}', $true, $false, $false)
$pres.SaveAs('{pdf}', 32)
$pres.Close()
$ppt.Quit()
"""
    try:
        res = subprocess.run(["powershell", "-NoProfile", "-Command", ps], capture_output=True, text=True, timeout=120)
        if res.returncode != 0:
            print(res.stderr)
    except Exception as exc:
        print(f"PowerPoint PDF export failed for {pptx}: {exc}")


def scrub_metadata() -> None:
    for path in MAIN.rglob("*.docx"):
        if path.name.startswith("~$"):
            continue
        try:
            doc = Document(path)
            set_doc_props(doc, doc.core_properties.title or "")
            style_doc(doc)
            doc.save(path)
        except Exception as exc:
            print(f"metadata docx skip {path}: {exc}")
    for path in MAIN.rglob("*.pptx"):
        if path.name.startswith("~$"):
            continue
        try:
            prs = Presentation(path)
            prs.core_properties.author = STUDENT
            prs.core_properties.last_modified_by = STUDENT
            prs.core_properties.comments = ""
            prs.save(path)
        except Exception as exc:
            print(f"metadata pptx skip {path}: {exc}")


def count_pdf_pages(pdf: Path) -> int | None:
    if not pdf.exists():
        return None
    try:
        return len(PdfReader(str(pdf)).pages)
    except Exception:
        return None


def update_readme() -> None:
    readme = MAIN / "README.md"
    text = f"""# {TITLE}

Nama Mahasiswa: {STUDENT}
NIM: {NIM}
Program Studi: {PROGRAM}
Fakultas: {FACULTY}
Universitas: {UNIVERSITY}
Tahun: {YEAR}

## Research Gap
Sebagian penelitian intrusion detection menekankan akurasi model, tetapi belum selalu menyediakan dashboard edukatif, klasifikasi risiko otomatis, rekomendasi tindakan, dan paket artefak yang mudah direproduksi mahasiswa.

## Novelty
Integrasi machine learning anomaly detection dengan dashboard monitoring ringan, risk score otomatis, dan recommended action untuk administrator jaringan pemula.

## Research Method
RM1: Membangun pipeline supervised learning untuk deteksi anomali jaringan menggunakan CICIDS2017.

RM2: Mengintegrasikan hasil model terbaik ke dashboard dan report export yang menampilkan risiko operasional.

## Ringkasan Hasil
Dataset: CICIDS2017 Friday Afternoon DDoS representative subset.

Model terbaik: Decision Tree.

Accuracy: 0.9992; Precision: 0.9995; Recall: 0.9990; F1-score: 0.9992.

## Struktur Folder
01_Paper sampai 15_Bukti_Submit mengikuti struktur artefak penelitian UAS AI, termasuk paper, mapping, gap analysis, dataset, source code, model, hasil eksperimen, visualisasi, draft IEEE, presentasi, Turnitin, deployment, GitHub, dokumentasi, dan bukti submit.
"""
    readme.write_text(text, encoding="utf-8")


def main() -> None:
    ensure_dirs()
    metrics, comparison, df = load_metrics_and_data()
    create_visuals(metrics, comparison, df)
    paper_index_docs()
    dataset_docs(metrics, df)
    experiment_docs(metrics, comparison)
    ieee_docx = build_ieee_doc(metrics, comparison)
    pptx = build_ppt(metrics, comparison)
    add_supporting_docs()
    update_readme()
    scrub_metadata()
    pages = export_docx_to_pdf(ieee_docx, MAIN / "09_Draft_IEEE" / "Draft_Artikel_IEEE.pdf")
    export_docx_to_pdf(MAIN / "01_Paper" / "Paper_Index.docx", MAIN / "01_Paper" / "Paper_Index.pdf")
    export_docx_to_pdf(MAIN / "04_Dataset" / "Dataset_Information.docx", MAIN / "04_Dataset" / "Dataset_Information.pdf")
    export_docx_to_pdf(MAIN / "07_Hasil_Eksperimen" / "Evaluation_Result_Detail.docx", MAIN / "07_Hasil_Eksperimen" / "Evaluation_Result_Detail.pdf")
    export_pptx_to_pdf(pptx, MAIN / "10_Presentasi" / "Slide_Presentasi.pdf")
    summary = {
        "ieee_docx": str(ieee_docx),
        "ieee_pdf_pages": pages or count_pdf_pages(MAIN / "09_Draft_IEEE" / "Draft_Artikel_IEEE.pdf"),
        "pptx": str(pptx),
        "ppt_pdf_pages": count_pdf_pages(MAIN / "10_Presentasi" / "Slide_Presentasi.pdf"),
    }
    (MAIN / "14_Dokumentasi" / "final_upgrade_summary.json").write_text(json.dumps(summary, indent=2), encoding="utf-8")
    print(json.dumps(summary, indent=2))


if __name__ == "__main__":
    main()
